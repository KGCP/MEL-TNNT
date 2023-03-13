'''
@component: Metadata Extraction & Loader (MEL).
@author: Sergio.
@summary: Metadata extraction and loading documents into CouchDB.
@project: AGRIF.
# History Update:
#    2019-10-04: internal release of version 0.1.0.
#    2020-01-17: recursive directory processing (sub-directories) + "Associated-Metadata" processing.
#    2020-04-08: correction of the text extraction method for PDF files.
#    2020-06-23: comprehensive CSV processing + encoding handling functions.
#    2020-06-25,26: NLP-NER task integration.
#    2020-08-07: OCR processing correction for PDF scanned documents.
#    2020-10-14,15: some generalizations for deployment in an OS-specific (Windows/Linux) environment.
#    2020-10-23,26: abstraction in the configuration file --> "dataset" concept
#    2020-11-09: improve the execution from the command line.
#    2020-11-17: logging.
#    2020-11-20: support for RTF (preserves table layouts).
#    2020-11-23~25: reorganising some code and improving interface.
#    2021-01-04: support for multiple configuration files.
#    2021-02-03: NHMRC_pruneMetadata() @ Test/CouchDB
#    2022-03-07: integration with the FutureSOILS Data Validation Engine (DVE).
#    2023-03-13: TNNT code improvements and bug correction.

# Specific dependencies:
http://www.decalage.info/python/oletools
https://textract.readthedocs.io/en/latest/installation.html
Linux server:
$ sudo yum update
$ sudo yum install tesseract

# OBSERVATIONS:

Sergio(2020-08-15): ** Microsoft Unicode Mappings **
If really necessary, we can implement the necessary character mappings from Microsfot specific encodings to Unicode.  See:
Source: <http://www.unicode.org/Public/MAPPINGS/VENDORS/MICSFT/>

Sergio(2020-11-18): ** REFACTORING? **
The main classes of the package have circular dependencies among them.
It's not worth it to decompose the whole script into various modules (one per class).
I tried to do it but the code started to look "messy" due to the inline *imports* as needed.
Besides, it seems that there were some unexpected results from the execution flow.

Sergio(2020-11-23): importing locally specific packages?
    def some_function():
        __fnc__ = sys._getframe().f_code.co_name + '#' + str(sys._getframe().f_code.co_firstlineno) # function + number of 1st line the source code
        if (__fnc__ not in Utils.IMPORTED_PACKAGES): # not in the array?
            import package
            from package import something
            Utils.IMPORTED_PACKAGES.append(__fnc__) # marked as already imported
'''


# ==================================================================================================
# Core libraries (packages):
import platform # Checking the O.S. | 'Windows' specific variables.
import requests # HTTP operations: (to CouchDB)...
import json
import os
import re
import sys
import datetime
import random
import logging
from logging import handlers
import builtins

# Specific packages/modules (local context) -- usage in their own respective functions:
# Utils.getClassNameFromFrame():
import inspect
# CouchDB.search():
from typing import Pattern
# Text.__init__():
import unicodedata
# Text.extractKeywords():
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
# File.setAttributes():
import mimetypes
import shutil
# File.extractMetadata():
import xml.etree.ElementTree as ET
from xmljson import parker
# File.extractPDFinfo():
from PyPDF4 import PdfFileReader
import subprocess # for PDF text extraction.
import textract # OCR ==> dependencies @ poppler-0.68.0/bin
# File.extractPPTXinfo():
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
# File.extractDOCXinfo(): # Python package <https://python-docx.readthedocs.io/en/latest/user/install.html> # *python-docx*
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from docx.oxml.text import run, paragraph
# File.extractOLEinfo():
import olefile
from olefile.olefile import OleMetadata
# File.extractOLEinfo(), File.extractMSGinfo():
if (platform.system() == 'Windows'):
    import win32com.client as win32 # for .msg, and .doc files
# File.extractCSVinfo():
import csv
from collections import defaultdict
# File.extractXLSXinfo():
import pyexcel as pe
# File.extractZIPinfo():
import zipfile
# File.extractRTFinfo():
from striprtf.striprtf import rtf_to_text
# File.extractContent():
import base64 # document content codification
# NER.output():
from hashlib import blake2s
# Utils.all_encodings()
import pkgutil
import encodings
# Utils.dict_update()
import collections.abc
#from dict_recursive_update import recursive_update # a different implementation


# ==================================================================================================
if hasattr(builtins, "MEL_config_filename"):
    _config_filename = builtins.MEL_config_filename
else:
    _config_filename = "config" # default
# check non-usage: default configuration filename:
if (_config_filename is None) or (not _config_filename):
    _config_filename = "config" # default

def init(cnfg_fn=_config_filename): # general initialization function
    Utils.loadConfigFile(cnfg_fn)
    Utils.initLogging() # logging


# ==================================================================================================
class Utils:
    global _config_filename
    CONFIG_JSON_DIR_ = os.path.dirname(os.path.realpath(__file__)) + "/"
        # "E:/Workspace/KG-Tools/MEL/" if (platform.system() == 'Windows') else\
        # "/data/mapping-services/MEL/"
    CONFIG_JSON_FILE = _config_filename + ".json" # default configuration file from the *builtins.MEL_config_filename* variable (if exists).
    if not(os.path.exists(CONFIG_JSON_FILE)):
        CONFIG_JSON_FILE = CONFIG_JSON_DIR_ + CONFIG_JSON_FILE  
    with open(CONFIG_JSON_FILE) as _config_json_f: # mode="rt" (default)
        _config = json.loads(_config_json_f.read())

    IMPORTED_PACKAGES = [] # experimental feature to keep track of imported packages
    
    DATASET_ID = str(_config["$-exec"]) # project (dataset in usage) string ID
    DATASET  = _config["Datasets"][DATASET_ID] # the local dataset configuration
    DATASETS = _config["Datasets"]
    
    _log = logging.getLogger() # No name assigned to the logger.
    _log_active = False
    
    _datetime2string_format = _config["DateTime-to-String-Format"] # Parameter for datetime.strftime(Utils._datetime2string_format)

    # Path file separator:
    FILE_PATH_SEPARATOR = "\\" if (platform.system() == 'Windows') else "/" # in Linux env.

    # Constants for Windows file path management:
    # https://docs.microsoft.com/en-us/windows/win32/fileio/naming-a-file?redirectedfrom=MSDN
    WinAPI_MAX_PATH = 260 if (platform.system() == 'Windows') else 5000
    WinAPI_ExtendedLengthPathPrefix = "\\\\?\\" if (platform.system() == 'Windows') else "" # Win: support for Unicode characters. | N/A in Linux
    
    # loads a specific configuration file.
    @staticmethod
    def loadConfigFile(_filename=_config_filename):
        if ( (_filename is not None) or (not _filename) ):
            if (Utils.CONFIG_JSON_FILE.endswith(str(_filename) + ".json") and (Utils._config is not None)):
                Utils.displayConfigFile("Already using", "nothing to do")
                return
        Utils.CONFIG_JSON_FILE  = "config" if ( (_filename is None) or (not _filename) or (_filename == "config") ) else str(_filename)
        Utils.CONFIG_JSON_FILE += ".json" # extension of the configuration file
        if not(os.path.exists(Utils.CONFIG_JSON_FILE)):
            Utils.CONFIG_JSON_FILE = Utils.CONFIG_JSON_DIR_ + Utils.CONFIG_JSON_FILE  
        with open(Utils.CONFIG_JSON_FILE) as _config_json_f: # mode="rt" (default)
            Utils._config = json.loads(_config_json_f.read())
            Utils.displayConfigFile("Set usage of", "loaded")
        Utils.DATASET_ID = str(Utils._config["$-exec"]) # project (dataset in usage) string ID
        Utils.DATASET  = Utils._config["Datasets"][Utils.DATASET_ID] # the local dataset configuration
        Utils.DATASETS = Utils._config["Datasets"]
        Dataset.loadDefault() # proceeds to load the default dataset of the configuration file.
    
    @staticmethod
    def displayConfigFile(msg, note):
        print(f"******************************************************************************")
        Utils.output(f"* {msg} configuration file: {Utils.CONFIG_JSON_FILE}\n  ({note})", _print=True)
        print(f"******************************************************************************\n")

    # Gets a datetime value with (possible) timezone ('tz') information in ISO 8601 format (YYYY-MM-DDTHH:MM:SS | YYYY-MM-DDTHH:MM:SS+HH:MM).
    # Reference: <https://docs.python.org/3.6/library/datetime.html#datetime.datetime>
    # IMPORTANT: The datetime values are formatted along with their time zone information (if available).
    # The program doesn't apply any kind of datetime value conversion among values that might be offset.
    @staticmethod
    def getFormattedDT(dt):
        Utils.output(f"dt={dt}; tzinfo={dt.tzinfo}", _print=False)
        ''' We don't apply any kind of datetime conversion:
        return dt.astimezone(tz=timezone.utc).isoformat(' ') if (dt.tzinfo) else ___ '''
        return dt.isoformat() # ISO 8601 format

    # Gets the folder path with the 32K extended-length for Win-API:
    @staticmethod
    def getFolder(folder=None): # folder is optional.
        """ * Allowed values: {"Input-Folder", "Temp-Folder", "Metadata-Folder", "Test-Folder"}
            * Default value: "Input-Folder" """
        return (Utils.WinAPI_ExtendedLengthPathPrefix + 
                Dataset._f.get(folder, 
                Dataset._f["Input-Folder"]))

    @staticmethod
    def getClassNameFromFrame(fr):
        return fr.f_locals.__class__.__name__ if (("self" in fr.f_locals) or ("cls" in fr.f_locals)) else "<static_method>"

    # General output to the console:
    @staticmethod
    def output(msg, _error=False, _print=bool(_config["Print-Verbose-Output"])):
        ERR = " [ERROR!]" if (_error) else ""
        prev_frame = inspect.currentframe().f_back
        fnc = Utils.getClassNameFromFrame(prev_frame) + "." + prev_frame.f_code.co_name # caller function.
        if (_print):
            if (msg.strip(" \t\n")):
                output_string = f"\n{datetime.datetime.now()}> @{fnc}{ERR}:\n{msg}"
                # code.interact(local=dict(globals(), **locals())) # for debugging...
                if (Utils._log_active):
                    Utils._log.info(output_string)
                else:
                    print(output_string)
            else:
                print(msg, end='')
        else:
            print("", end='')

    # Start time stamp:
    @staticmethod
    def printStartTimeStamp(dt_begin):
        Utils._log.info(f"Start time: [{dt_begin}]")

    # End time stamp:
    @staticmethod
    def printEndTimeStamp(dt_end, delta):
        Utils._log.info(f"\n\nEnd time: [{dt_end}]")
        Utils._log.info(f"Execution time: {delta} seconds.")

    # Get a boolean value from a dictionary[key] element.  Checks if the [key] exists; if not, returns False.
    @staticmethod
    def getsBoolFromDictKey(d: dict, k: str) -> bool:
        return (bool(d[k]) if (k in d) else False)

    @staticmethod
    def addItemInCountingList(_list, key):
        if (key in _list.keys()):
            _list[key] += 1 # counting repetitions
        else: # first time the element is added in the list
            _list[key] = 1

    @staticmethod
    def all_encodings():
        modnames = set([modname for importer, modname, ispkg in pkgutil.walk_packages(
            path=[os.path.dirname(encodings.__file__)], prefix='')])
        aliases = set(encodings.aliases.aliases.values())
        return modnames.union(aliases)

    @staticmethod
    def getValidEncoding(text):
        for enc in Utils.all_encodings():
            try:
                msg = text.decode(enc)
            except UnicodeEncodeError:
                continue
            else:
                if len(msg) > 0:
                    Utils._log.info(f'Decoding {text} with {enc} is {msg}')
    
    @staticmethod
    def isJSON(_str):
        try:
            JSON_object = json.loads(str(_str))
        except ValueError as _:
            return False
        return True

    @staticmethod
    def tryPrettyfiedJSON(obj, _indent=4):
        _prettyfiedJSON = ""
        try:
            _prettyfiedJSON = json.dumps(obj, indent=_indent)
        except Exception as _:
            return str(obj)
        return _prettyfiedJSON

    @staticmethod
    def dict_update(d, u):
        for k, v in u.items():
            if isinstance(v, collections.abc.Mapping):
                d[k] = Utils.dict_update(d.get(k, {}), v)
            else:
                d[k] = v
        return d
    
    @staticmethod
    def initLogging(logfile=""):
        if (not Utils._log_active):
            logging.root.handlers = [] # reset any system configuration about the logging.
            Utils._log.setLevel(logging.INFO)
            file_format = logging.Formatter(fmt=Utils._config["Logging"]["Format"])
            # to the console:
            ch = logging.StreamHandler(sys.stdout)
            ch.setLevel(logging.INFO)
            ch.setFormatter("") # no format
            Utils._log.addHandler(ch)
            # to a file:
            fh = logging.FileHandler(\
                ( (Utils._config["Logging"]["Default-Folder"] + logfile + '.log') if (logfile) else Dataset._log ),\
                mode='w', encoding='utf-8')
            fh.setLevel(logging.INFO)
            fh.setFormatter(file_format)
            Utils._log.addHandler(fh)
            Utils._log_active = True

Utils.displayConfigFile("Default", "loaded")


# ==================================================================================================
class Dataset: # Local Dataset Configuration
    _id = Utils.DATASET_ID
    _   = Utils.DATASET
    _db  = _["CouchDB"]
    _f   = _["Folders"]
    _ner = _["NLP-NER"]
    _log = _f["Log-Folder"] + _id + '.log'
    '''
    Output File Generation: COMPLETE (metadata + NER results), ONLY-NER
    '''
    
    @staticmethod
    def load(_id):
        if (_id not in Utils.DATASETS):
            Utils.output(f"Dataset *{_id}* not found in the configuration file!\nThe current dataset ({Dataset._id}) will be kept.", _error=True, _print=True)
            return
        if (_id == Dataset._id): # prevents loading the current dataset settings 
            return
        global DB
        Dataset._id = _id
        Dataset._   = Utils.DATASETS[Dataset._id]
        Dataset._db  = Dataset._["CouchDB"]
        DB = Dataset._db["Database"]
        Dataset._f   = Dataset._["Folders"]
        Dataset._ner = Dataset._["NLP-NER"]
        Dataset._log = Dataset._f["Log-Folder"] + Dataset._id + '.log'
        # Update the CouchDB-related variables:
        CouchDB._end_point     = CouchDB.getEndPoint(Dataset._db ["Database"])
        CouchDB._NER_end_point = CouchDB.getEndPoint(Dataset._ner["Database"])
    
    @staticmethod
    def loadDefault():
        Dataset.load(Utils.DATASET_ID)
    
    @staticmethod
    def couchDBenabled():
        return ( bool(CouchDB._couchDB["Enable"]) and bool(Dataset._db["Enable"]) )
    
    @staticmethod
    def NERenabled():
        return ( bool(NER.config["Enable"]) and bool(Dataset._ner["Enable"]) )

    @staticmethod
    def NER_Complete():
        return ( Dataset.NERenabled() and (Dataset._ner["Output-Handling"] == "COMPLETE") )
    
    @staticmethod
    def NER_OnlyNER():
        return ( Dataset.NERenabled() and (Dataset._ner["Output-Handling"] == "ONLY-NER") )

    @staticmethod
    def NER_storeOutputOnCouchDB():
        return ( Dataset.NERenabled() and bool(Dataset._ner["Store-Output-On-CouchDB"]) )

    @staticmethod
    def isGenerateOutputFileEnabled() -> bool:
        return Utils.getsBoolFromDictKey(Dataset._f, "Generate-Output-Files")

    @staticmethod
    def hasAssociatedMetadata() -> bool:
        return Utils.getsBoolFromDictKey(Dataset._f, "Has-Associated-Metadata")

    @staticmethod
    def isViewFilePrintEnabled() -> bool:
        return Utils.getsBoolFromDictKey(Dataset._f, "View-File-Print")

    @staticmethod
    def pdftotext_preserveLayout() -> bool:
        return Utils.getsBoolFromDictKey(Dataset._f, "PDF.Text-Extraction-Preserve-Layout")
    
    @staticmethod
    def writeOutputFile(_filename, sufix, output):
        if (Dataset.isGenerateOutputFileEnabled()):
            output_filename = Dataset._f["Output-Folder"] + _filename + sufix 
            with open(output_filename, 'w') as outfile:
                outfile.write( Utils.tryPrettyfiedJSON(output) )
            Utils.output(f"[{output_filename}]", _print=True)
    @staticmethod
    def generateMergedOutput(_filename, output): # (string, json) # we assume that *output* has the full content including the NER results...
        if (Dataset.NER_Complete()):
            Dataset.writeOutputFile(_filename, "-MEL+NER_output.json", output)
    @staticmethod
    def generateNERoutput(_filename, output): # (string, json) # we assume that *output* has only the NER results...
        if (Dataset.NER_OnlyNER()):
            Dataset.writeOutputFile(_filename, "-NER_output.json", output)

    @staticmethod
    def generateOutputFile(_filename, obj, processNER=False):
        if (Dataset.isGenerateOutputFileEnabled()):
            output, NER_result = {}, {}
            if (isinstance(obj, File)):
                output = obj.extractMetadata()  # if (COMPLETE): it's already processed inside the metadata extraction
                if (Dataset.NER_Complete()):
                    Dataset.generateMergedOutput(_filename, output)
                    return
                processNER = Dataset.NER_OnlyNER() # NER processing?
            if (isinstance(obj, dict)):
                output = obj
            # NER processing or retrieving the last NER result:
            NER_result = NER.process(output) if (processNER) else NER._last_result
            if   (Dataset.NER_OnlyNER()):
                Dataset.generateNERoutput(_filename, NER_result)
            elif (Dataset.NER_Complete()):
                output.update(NER_result)
                Dataset.generateMergedOutput(_filename, output)
            else:
                Dataset.writeOutputFile(_filename, "-MEL_output.json", output)

    @staticmethod
    def storeOutputOnCouchDB():
        return ( (Dataset.NERenabled()) and bool(Dataset._ner["Store-Output-On-CouchDB"]) )
    
    @staticmethod
    def reloadFile(extension):
        return ( bool(CouchDB._couchDB["Operations"]["Reload-Documents"]) and bool(Dataset._db["Reload-Documents"]) ) or\
                (extension.upper() in Dataset._db["Reload-Files-of-Extensions"].strip().upper())

    @staticmethod
    def NER_regenerateOutput():
        return ( Dataset.NERenabled() and bool(Dataset._ner["Regenerate-Output"]) )

    @staticmethod
    def shortPrint():
        Utils._log.info(f"CouchDB Processing: [{bool(Dataset._db['Enable'])}], NLP-NER Processing: [{bool(Dataset._ner['Enable'])}]")
        Utils._log.info(f"Input Folder: [{Utils.getFolder()}]\n\n") # "Input-Folder"

# default/current database:
DB = Dataset._db["Database"]


# ==================================================================================================
class CouchDB:
    _couchDB = Utils._config["CouchDB"]
    _host = _couchDB["Host"]
    _end_point     = _host + Dataset._db ["Database"]
    _NER_end_point = _host + Dataset._ner["Database"]

    @staticmethod
    def getEndPoint(db):
        return (CouchDB._host + db)

    @staticmethod
    def query(db, query={}):
        Utils.output(json.dumps(query, indent=4), _print=False)
        req = requests.post(CouchDB._host + db + CouchDB._couchDB["Operations"]["Find"], json=query)
        # //result = json.loads(req.text)
        # //Utils.output(json.dumps(result, indent=4), _print=False)
        json_response = json.loads(req.content)
        if (req.status_code == 200): # 200 OK - Request completed successfully
            Utils.output(f"""*** {json_response['execution_stats']['results_returned']} were found *** | Details below: {
                                  json.dumps(json_response, indent=4)}""", _print=False)
            return json_response
        else:
            # Possible HTTP Response:
            # Status Codes:    
            #    400 Bad Request - Invalid request
            #    401 Unauthorized - Read permission required
            #    500 Internal Server Error - Query execution error
            Utils.output(f"""The query execution failed:
                * HTTP Response Code = {req.status_code}
                * HTTP Headers = {req.headers}""", _error=True, _print=True)
            Utils.output(json.dumps(json_response, indent=4), _error=True, _print=True)
            return req

    @staticmethod
    def queryDocs(db, selector={}, only_count=True, fields=["_id", "_rev"], limit=10000, skip=0):
        query = {
            "selector": selector,
            "fields": fields,
            "limit": limit,
            "skip": skip,
            "execution_stats": True
        }
        Utils.output(json.dumps(query, indent=4), _print=False)
        r = CouchDB.query(db, query)
        return r['execution_stats']['results_returned'] if (only_count) else r

    @staticmethod
    def getDocumentSet(_db=DB, _re_path="", _re_filename="", _list_extensions=None, _limit=10, _chk_text=True):
        selector = { "General-Metadata": { } }
        _limit = 10 if (_limit <= 0) else _limit # default value
        def _re(attribute, re_pattern):
            if (re_pattern):
                re_pattern = "(.*)(" + re.sub(r'([\^\$\.\+\(\)\[\]\{\}\\])', r"\\\1", re_pattern) + ")(.*)"
                selector["General-Metadata"].update( { attribute : { "$regex": re_pattern } } )
        _re("ABSOLUTEPATH", _re_path)
        _re("FILENAME"    , _re_filename)
        # (_list_extensions is None) ==> look for all file extensions
        _list_extensions = _list_extensions if (_list_extensions is not None) else []
        _or_extensions = []
        for extension in _list_extensions:
            _or_extensions.append({ "EXTENSION": extension.lower() })
            _or_extensions.append({ "EXTENSION": extension.upper() })
        if (_or_extensions):
            selector["General-Metadata"].update( { "$or": _or_extensions } )
        fields = ["_id", "General-Metadata", "Use-Case$Folder", "Specific-Metadata"]
        Utils._log.info(f"INPUT PARAMETERS:\n* Dataset={_db}\n* Selector:\n{json.dumps(selector, indent=4)}\n* Fields={fields}\n* Limit={_limit}\n")
        r = CouchDB.queryDocs(_db, selector, False, fields, _limit)
        Utils.output(f"db={{{_db}}}:\n{json.dumps(r, indent=4)}", _print=False)
        doc_set = []
        n = 1
        text = ""
        for doc in r["docs"]:
            # Document's text content:
            if (_chk_text):
                text = doc["Specific-Metadata"]["text-analysis"]["clean-text"] \
                    if  ("Specific-Metadata" in doc) and \
                        ("text-analysis" in doc["Specific-Metadata"]) and \
                        ("clean-text" in doc["Specific-Metadata"]["text-analysis"]) \
                    else ""
                if (len(text) == 0):
                    continue # next doc.
            Utils.output(f"""#{n}: id={doc["_id"]}\nf={doc["General-Metadata"]["ABSOLUTEPATH"]}""", _print=False)
            n += 1
            doc_set.append(doc) # add document to the result
        Utils._log.info(f"""Documents found in *CouchDB*: {len(r["docs"])}\nDocuments with text content: {len(doc_set)}""")
        return doc_set
    
    @staticmethod
    def applyREresults(_re_pattern, _results):
        return {
            "$PATTERN$": _re_pattern if (isinstance(_re_pattern, str)) else _re_pattern.pattern, # RE pattern
            "$RESULT$" : _results
        }
    
    @staticmethod
    def getCompletePattern(_re_txt_search_BEGIN_pattern="", _re_txt_search_END_pattern=""):
        if (_re_txt_search_BEGIN_pattern): # non-capturing version of regular parentheses.
            _re_txt_search_BEGIN_pattern = "(?:" + _re_txt_search_BEGIN_pattern + ")"
        _re_txt_search_CHUNK_pattern = "(.*)"
        if (_re_txt_search_END_pattern): # non-capturing version of regular parentheses.
            _re_txt_search_CHUNK_pattern += "(?:" + _re_txt_search_END_pattern + ")"
            _re_txt_search_END_pattern = ""
        _re_complete_pattern = _re_txt_search_BEGIN_pattern + _re_txt_search_CHUNK_pattern + _re_txt_search_END_pattern
        if (_re_complete_pattern == "(.*)"): # the complete pattern will match with the full text content
            Utils.output(f"The search pattern will retrieve the full text content: {_re_complete_pattern}", _error=True, _print=True)
            return _re_complete_pattern, None
        _re_txt_search_pattern = re.compile(_re_complete_pattern, flags=0)
        Utils.output(f"Search pattern: {_re_txt_search_pattern}\n", _print=True)
        return _re_complete_pattern, _re_txt_search_pattern

    @staticmethod
    def getFirstChunk(_db=DB, _re_path="", _re_filename="", _list_extensions=None, _limit=10, \
                      _re_txt_search_BEGIN_pattern="", _re_txt_search_END_pattern=""):
        _re_complete_pattern, _re_txt_search_pattern = CouchDB.getCompletePattern(_re_txt_search_BEGIN_pattern, _re_txt_search_END_pattern)
        if (_re_complete_pattern == "(.*)"): # the complete pattern will match with the full text content
            return []
        # The text search is performed in the Python script: we avoid using CouchDB's (document store) searching mechanism
        # due that the database is not indexed, and applying (complex) regular expression in large contents could detriment the response time.
        docs = CouchDB.getDocumentSet(_db=DB, _re_path=_re_path, _re_filename=_re_filename, _list_extensions=_list_extensions, _limit=_limit)
        doc_set = []
        n = 1
        text, match = "", None
        for doc in docs:
            match = None
            text = doc["Specific-Metadata"]["text-analysis"]["clean-text"] # Document's text content (we know that this structure exists)
            match = _re_txt_search_pattern.search(text) # Returns the first match
            if not(match):
                continue # next doc.
            doc.update( CouchDB.applyREresults(_re_complete_pattern, match) )
            Utils.output(f"""#{n}: id={doc["_id"]}\nf={doc["General-Metadata"]["ABSOLUTEPATH"]}""", _print=True)
            if (match is not None):
                Utils.output(f"m={match}\n\n", _print=True)
            else:
                Utils.output("\n", _print=True)
            n += 1
            doc_set.append(doc) # add document (with the result -- matching object)
        Utils.output(f"""Document set cardinality: {len(docs)}\nChunk capturing matches: {len(doc_set)}""", _print=True)
        return doc_set

    @staticmethod
    def getAllChunks(_db=DB, _re_path="", _re_filename="", _list_extensions=None, _limit=10, \
                      _re_txt_search_BEGIN_pattern="", _re_txt_search_END_pattern=""):
        _re_complete_pattern, _re_txt_search_pattern = CouchDB.getCompletePattern(_re_txt_search_BEGIN_pattern, _re_txt_search_END_pattern)
        if (_re_complete_pattern == "(.*)"): # the complete pattern will match with the full text content
            return []
        # The text search is performed in the Python script: we avoid using CouchDB's (document store) searching mechanism
        # due that the database is not indexed, and applying (complex) regular expression in large contents could detriment the response time.
        docs = CouchDB.getDocumentSet(_db=DB, _re_path=_re_path, _re_filename=_re_filename, _list_extensions=_list_extensions, _limit=_limit)
        doc_set = []
        n = 1
        text, matches = "", []
        for doc in docs:
            matches = []
            text = doc["Specific-Metadata"]["text-analysis"]["clean-text"] # Document's text content (we know that this structure exists)
            for m in _re_txt_search_pattern.finditer(text):
                if (m):
                    matches.append(m)
            if not(matches):
                continue # next doc.
            doc.update( CouchDB.applyREresults(_re_complete_pattern, matches) )
            Utils.output(f"""#{n}: id={doc["_id"]}\nf={doc["General-Metadata"]["ABSOLUTEPATH"]}""", _print=True)
            if (matches):
                Utils.output(f"m={matches}\n\n", _print=True)
            else:
                Utils.output("\n", _print=True)
            n += 1
            doc_set.append(doc) # add document (with the result -- matches array)
        Utils._log.info(f"""Document set cardinality: {len(docs)}\nText search matches: {len(doc_set)}""")
        return doc_set

    @staticmethod
    def searchFirstMatch(_db=DB, _re_path="", _re_filename="", _list_extensions=None, _limit=10, _re_txt_search_pattern=None):
        # The text search is performed in the Python script: we avoid using CouchDB's (document store) searching mechanism
        # due that the database is not indexed, and applying (complex) regular expression in large contents could detriment the response time.
        docs = CouchDB.getDocumentSet(_db=DB, _re_path=_re_path, _re_filename=_re_filename, _list_extensions=_list_extensions, _limit=_limit)
        doc_set = []
        n = 1
        text, match = "", None
        for doc in docs:
            match = None
            text = doc["Specific-Metadata"]["text-analysis"]["clean-text"] # Document's text content (we know that this structure exists)
            if (_re_txt_search_pattern is not None) and isinstance(_re_txt_search_pattern, Pattern):
                # It's expected that *_re_txt_search_pattern* is an instance of type "Pattern":
                # _re_txt_search_pattern = re.compile(pattern, flags=0) # valid
                match = _re_txt_search_pattern.search(text) # Returns the first match
                if not(match):
                    continue # next doc.
                doc.update( CouchDB.applyREresults(_re_txt_search_pattern, match) )
            Utils.output(f"""#{n}: id={doc["_id"]}\nf={doc["General-Metadata"]["ABSOLUTEPATH"]}""", _print=False)
            if (match is not None):
                Utils.output(f"m={match}\n\n", _print=False)
            else:
                Utils.output("\n", _print=False)
            n += 1
            doc_set.append(doc) # add document (with the result -- matching object)
        Utils._log.info(f"""Document set cardinality: {len(docs)}\nText search matches: {len(doc_set)}""")
        return doc_set

    @staticmethod
    def searchAllMatches(_db=DB, _re_path="", _re_filename="", _list_extensions=None, _limit=10, _re_txt_search_pattern=None):
        # The text search is performed in the Python script: we avoid using CouchDB's (document store) searching mechanism
        # due that the database is not indexed, and applying (complex) regular expression in large contents could detriment the response time.
        docs = CouchDB.getDocumentSet(_db=DB, _re_path=_re_path, _re_filename=_re_filename, _list_extensions=_list_extensions, _limit=_limit)
        doc_set = []
        n = 1
        text, matches = "", []
        for doc in docs:
            matches = []
            text = doc["Specific-Metadata"]["text-analysis"]["clean-text"] # Document's text content (we know that this structure exists)
            if (_re_txt_search_pattern is not None) and isinstance(_re_txt_search_pattern, Pattern):
                # It's expected that *_re_txt_search_pattern* is an instance of type "Pattern":
                # _re_txt_search_pattern = re.compile(pattern, flags=0) # valid
                for m in _re_txt_search_pattern.finditer(text):
                    if (m):
                        matches.append(m)
                if not(matches):
                    continue # next doc.
                doc.update( CouchDB.applyREresults(_re_txt_search_pattern, matches) )
            Utils.output(f"""#{n}: id={doc["_id"]}\nf={doc["General-Metadata"]["ABSOLUTEPATH"]}""", _print=True)
            if (matches):
                Utils.output(f"m={matches}\n\n", _print=True)
            else:
                Utils.output("\n", _print=True)
            n += 1
            doc_set.append(doc) # add document (with the result -- matching object)
        Utils._log.info(f"""Document set cardinality: {len(docs)}\nText search matches: {len(doc_set)}""")
        return doc_set

    @staticmethod
    def checkDuplicateDocs(db):
        Utils.output(f"*** Checking for duplicated documents in the database {db} ***", _print=True)
        selector = { }
        fields = ["_id", "General-Metadata"]
        r = CouchDB.queryDocs(db, selector, False, fields)
        n, m = 1, 0
        duplicated = set()
        for x in r["docs"]:
            if not (x["_id"] in duplicated):
                processed_x = False
                for y in r["docs"]:
                    if not (y["_id"] in duplicated):
                        if ((x["General-Metadata"]["ABSOLUTEPATH"] == y["General-Metadata"]["ABSOLUTEPATH"]) and
                            x["_id"] != y["_id"]): # x and y are duplicated!
                            if not (processed_x):
                                m += 1 # duplicated docs: add x.
                                Utils._log.info(f"""\nDocument #{n}: {{{x["General-Metadata"]["ABSOLUTEPATH"]}}}""")
                                Utils._log.info(f"""id={x["_id"]}; (duplicated counter: {m})""")
                                duplicated.add(x["_id"])
                                processed_x = True
                            m += 1 # duplicated docs: add y.   
                            Utils._log.info(f"""id={y["_id"]}; (duplicated counter: {m})""")
                            duplicated.add(y["_id"])
            n += 1
        Utils.output(f"\nDuplicated docs: {m}; Revised docs: {len(r['docs'])}")

    @staticmethod
    def countDocsXsubPath(db, subpath):
        selector = {
            "General-Metadata": {
                "ABSOLUTEPATH": {
                "$regex": "(" + re.sub(r'([\?\\])', r"\\\1", Utils.WinAPI_ExtendedLengthPathPrefix) + ")?(.*)(" + subpath + ")(.*)"
                }
            }
        }
        return CouchDB.queryDocs(db, selector)

    @staticmethod
    def findDocument(file, db=Dataset._db["Database"]): # file is an instance of the *File* class
        if not(isinstance(file, File)):
            Utils.output("Expected an instance of class *File*.", _error=True, _print=True)
            return {}
        '''
        (1) The search is done via the field *ABSOLUTEPATH*, which we know that it always appear.
            ** There's a special case for the content of ZIP files and e-mail attachments of MSG files.
        (2) The following query was modified to include a regular expression as part of the search to look for documents that were
        loaded before with regular Windows paths along with documents that were loaded with extended-length Windows paths.
        The lines of code include the handling of the regex metacharacters as part of the file path. '''
        query = {
            "selector": {
                "General-Metadata": {
                    "FILENAME": file.name,
                    "ABSOLUTEPATH": {
                        "$regex": "(" + re.sub(r'([\\\?])', r"\\\1", file.path[:3]) + ")?" +
                                        re.sub(r'([\^\$\.\+\(\)\[\]\{\}\\])', r"\\\1", file.path[4:])
                    }
                },
                "Use-Case$Folder": file.useCase
            },
            "fields": ["_id", "_rev", "General-Metadata", "Use-Case$Folder"],
            "limit": 100,
            "skip": 0,
            "execution_stats": True
        }
        Utils.output(json.dumps(query, indent=4), _print=False)
        return CouchDB.query(db, query)

    @staticmethod
    def getDocumentGivenID(_id, db=Dataset._db["Database"]):
        end_point = CouchDB._host + db + '/' + _id
        Utils.output(f"""* CouchDB endpoint = [{end_point}]""", _print=False)
        req = requests.get(end_point)
        # HTTP Response:
        json_response = json.loads(req.content)
        Utils.output(f"""
            * HTTP Response Code = {req.status_code}
            * HTTP Headers = {req.headers}
            * JSON Response = {json.dumps(json_response, indent=4)}""", _print=False)
        '''
        if (req.status_code == 200):
            Utils._log.info(f"title: {json_response['title']}")
            Utils._log.info("_attachments: " + json.dumps(json_response["_attachments"], indent=4))
        '''
        return json_response

    @staticmethod
    def getDocumentGivenFilename(filename, db=Dataset._db["Database"]): # search based on *filename*
        re = filename.replace('.', '\\.').replace('(', '\\(').replace(')', '\\)')
        selector = {
            "General-Metadata": {
                "FILENAME": {
                    "$regex": f"({re})$"
                }
            }
        }
        fields = ["_id", "General-Metadata", "Specific-Metadata"]
        r = CouchDB.queryDocs(db, selector, False, fields)
        Utils.output(f"db={{{db}}}:\n{json.dumps(r, indent=4)}", _print=False)
        return r

    @staticmethod
    def deleteDocument(_id, _rev, _end_point=""): # identifier and revision number on CouchDB
        _end_point = CouchDB._end_point if (_end_point == "") else _end_point
        req = requests.delete(_end_point + '/' + _id + '?rev=' + _rev)
        # HTTP Response:
        json_response = json.loads(req.content)
        Utils.output(f"""
            * HTTP Response Code = {req.status_code}
            * HTTP Headers = {req.headers}
            * JSON Response = {json.dumps(json_response, indent=4)}""", _print=True)
        return ((req.status_code == 200) or (req.status_code == 201))

    @staticmethod
    def updateDocument(_id, payload, _end_point=""):
        _end_point = CouchDB._end_point if (_end_point == "") else _end_point
        # the revision number (_rev) is specified in the JSON object structure with the full-content to update:
        r = requests.put(_end_point + '/' + _id, json=payload)
        result = json.loads(r.text)
        Utils.output(json.dumps(result, indent=4), _print=False)
        return result

    @staticmethod
    def addDocument(file, _end_point=""): # file of type 'File'
        _end_point = CouchDB._end_point if (_end_point == "") else _end_point
        if not(isinstance(file, File)):
            Utils.output("Expected an instance of class *File*.", _error=True, _print=True)
            return {}
        metadata = file.extractMetadata() # returns a JSON object # if (COMPLETE), the JSON object includes the NER results... 
        if (file.occurredCriticalException()):
            Utils.output(f"A critical exception occurred during the extraction.\nThe document won't be added to CouchDB.", _error=True, _print=True)
            return None, None, metadata
        payload = metadata
        file_content = file.extractContent() # returns a JSON object
        if (file_content): # if not empty
            payload["_attachments"] = file_content
        Utils.output(json.dumps(payload, indent=4), _print=False)
        r = requests.post(_end_point, json=payload)
        result = json.loads(r.text)
        Utils.output(f"{r}\n" + json.dumps(result, indent=4), _print=False)
        return r, result, metadata # return the result and the extracted metadata


# ==================================================================================================
class AssociatedMetadata:
    # load(): Loads a specific (.xlsx) file with the associated metadata.
    @staticmethod
    def load(key):
        _settings = Utils._config["Associated-Metadata"][key + "-Settings"]
        _d = _settings["Path"] + _settings["Folder"]
        Directory.scanAndBuildStructure(_d)
        _f = File(path=(_d + Utils.FILE_PATH_SEPARATOR + _settings["File"]), useCase=_settings["Folder"],
                  defaultDirAttributes=Directory.structure[_settings["File"]], JSONmetadata={})
        DATA = _f.extractMetadata(_perform_NER_processing=False) # do not process NLP-NER for the associated metadata file
        Utils.output(json.dumps(DATA, indent=4), _print=False)
        return DATA


    @staticmethod
    def cast(d, value):
        if (d == "str"):
            return str(value).strip()
        elif (d == "int"):
            return int(value)
        elif (d == "float"):
            return float(value)
        else:
            return value # no casting!


    # Moves back the found position.
    # It will point to the first element of the list.
    @staticmethod
    def movePositionToFirstElement(a, x, d, i, p):
        if (p > 1): # lower bound = 1; row 0 is the spreadsheet header.
            while (AssociatedMetadata.cast(d, a[p-1][i]) == AssociatedMetadata.cast(d, x)):
                p -= 1
                if (p == 1):
                    break
        Utils.output(f"x={x}; a[{p}][{i}]={a[p][i]}; d='{d}'; p={p};", _print=False)
        return p # position of the found element (FIRST of the list)


    # simpleSearch(a, x, i, lo, hi):
    #     a: data array;  x: element (value) to search for;
    #     d: data type of the values to compare;
    #     i: index of the data field to look for "x";
    #     lo: lower bound; hi: higher bound.
    #     move: "++" or "--".
    # Performs a simple search on the loaded data based on the "move" ("++" or "--").
    @staticmethod
    def simpleSearch(a, x, d, i, lo, hi, move):
        p = lo if (move == "++") else hi
        while True:
            if   ((move == "++") and (AssociatedMetadata.cast(d, a[p][i]) == AssociatedMetadata.cast(d, x))):
                return p # found!
            elif ((move == "--") and (AssociatedMetadata.cast(d, a[p][i]) == AssociatedMetadata.cast(d, x))):
                return AssociatedMetadata.movePositionToFirstElement(a, x, d, i, p)
            p = p + (1 if (move == "++") else -1) # inc or dec
            if  ((move == "++") and (p > hi)) or\
                ((move == "--") and (p < lo)):
                break
        return -1 # not found


    # binarySearch(a, x, d, i, lo, hi):
    #     a: data array;  x: element (value) to search for;
    #     d: data type of the values to compare;
    #     i: index of the data field where the list is sorted;
    #     lo: lower bound; hi: higher bound.
    # Performs a binary search on the loaded data.
    # It's expected to be found multiple rows with the same key ("x").
    # The algorithm looks for the first repetition.
    @staticmethod
    def binarySearch(a, x, d, i, lo, hi):
        if (lo > hi): # the element was not found!
            if (lo >= len(a)): # the lower bound exceeds the maximun length.
                return -1 # not found!
            # Tries to perform a simple search only for string values (lexicographic sorting problem with Excel).
            if (d == "str"):
                if (len(x) > len(a[hi][i])):
                    return AssociatedMetadata.simpleSearch(a, x, d, i, hi, len(a)-1, "++")
                if (len(x) < len(a[lo][i])):
                    return AssociatedMetadata.simpleSearch(a, x, d, i, 1, lo, "--")
            return -1 # not found!
        mid = (lo+hi) // 2 # integer division
        Utils.output(f"x={x}; a[{mid}][{i}]={a[mid][i]}; d='{d}'; lo={lo}; hi={hi}", _print=False)
        if (AssociatedMetadata.cast(d, a[mid][i]) == AssociatedMetadata.cast(d, x)):
            return AssociatedMetadata.movePositionToFirstElement(a, x, d, i, mid)
        elif (AssociatedMetadata.cast(d, x) < AssociatedMetadata.cast(d, a[mid][i])): # the element is smaller than the mid point:
            return AssociatedMetadata.binarySearch(a, x, d, i, lo, mid-1)
        else: # the element is greater than the mid point:
            return AssociatedMetadata.binarySearch(a, x, d, i, mid+1, hi)


# ==================================================================================================
class Text:
    __slots__ = ('_text', '_ascii_text', '_cleanText', '_raw_text')

    def __init__(self, text=""):
        self._text = str(text) # text might be an array of bytes from PDF extraction, which is not JSON serializable.
        # from Unicode (array of bytes) to ASCII (string):
        self._ascii_text = unicodedata.normalize('NFKD', self._text).encode('ascii','ignore').decode("utf-8")
        self._cleanText = self.cleanText()
        # self._raw_text = text

    # This function will remove all stop words and punctuations in the text and return a list of keywords.
    def extractKeywords(self, _ascii=True):
        wordTokens = word_tokenize(self._ascii_text if _ascii else self._text) # Split the text words into tokens
        punctuations = ['(',')',';',':','[',']',','] # Remove blow punctuation in the list.
        stopWords = stopwords.words('english') # Get all stop words in English.
        keywords = {}
        for word in wordTokens:
            # Makes sure that the word matches alphanumeric characters at the beginning and end of the string.
            # The list will return only keywords that are not in stop words and punctuations:
            if (re.search(r'(^\w+$)|(\w+$)', word)) and (word not in stopWords) and (word not in punctuations):
                Utils.addItemInCountingList(keywords, word)
        return keywords
    
    def cleanText(self):
        # Cleaning-up the text before => Unicode replacements:
        # \u2013: '-'
        # \u2019: 'RIGHT SINGLE QUOTATION MARK' to "'"
        # \u00a0: 'NO-BREAK SPACE' to ' '
        # \u0007: 'BELL' / <Control> to ' '
        return self._ascii_text\
            .replace(u"\u2013", "-")\
            .replace(u"\u2019", "'")\
            .replace(u"\u00a0", " ")\
            .replace(u"\u0007", " ")
    
    def applyPatternMatching(self):
        jsonRegExps = {}
        # The following is a regular expression for the full syntax of e-mails:
        # (r"(?:[a-z0-9!#$%&'*+/=?^_`{|}~-]+(?:\.[a-z0-9!#$%&'*+/=?^_`{|}~-]+)*|""(?:[\x01-\x08\x0b\x0c\x0e-\x1f\x21\x23-\x5b\x5d-\x7f]|\\[\x01-\x09\x0b\x0c\x0e-\x7f])*"")@(?:(?:[a-z0-9](?:[a-z0-9-]*[a-z0-9])?\.)+[a-z0-9](?:[a-z0-9-]*[a-z0-9])?|\[(?:(?:(2(5[0-5]|[0-4][0-9])|1[0-9][0-9]|[1-9]?[0-9]))\.){3}(?:(2(5[0-5]|[0-4][0-9])|1[0-9][0-9]|[1-9]?[0-9])|[a-z0-9-]*[a-z0-9]:(?:[\x01-\x08\x0b\x0c\x0e-\x1f\x21-\x5a\x53-\x7f]|\\[\x01-\x09\x0b\x0c\x0e-\x7f])+)\])", text)
        for name, pattern in Utils._config["Regular-Expressions"].items():
            matches = {}
            prog = re.compile(pattern, flags=(re.ASCII|re.MULTILINE))
            _list = prog.findall(self._cleanText)
            for item in _list:
                # item is not string when the pattern has group definitions (item is a tuple of the pattern groups).
                key = "".join(item) if not(isinstance(item, str)) else item
                Utils.addItemInCountingList(matches, key)
            jsonRegExps[name] = matches
        Utils.output(json.dumps(jsonRegExps, indent=4), _print=False)
        return jsonRegExps
    
    def analysis(self):
        ana = {}
        ana["text"] = self._text
        ana["clean-text"] = self._cleanText
        ana["ascii-text"] = self._ascii_text
        # ana["raw-text"] = self._raw_text # the raw text could be an array of bytes which is not JSON serializable!
        ana["pattern-matching"] = self.applyPatternMatching()
        ana["number-of-characters"] = str(self._text.__len__())
        keywords = self.extractKeywords()
        ana["Tokenized-text$Keywords"] = keywords
        ana["Number-of-keywords"] = str(keywords.__len__())
        return ana


"""
# ==================================================================================================
# /* from html.parser import HTMLParser */

class _HTMLparser(HTMLParser):
    __slots__ = ('tag')

    def __init__(self):
        self.tag = ""
    
    def handle_starttag(self, tag, attrs):
        self.tag = tag
    
    def handle_data(self, data):
        if (self.tag == "title"):
            self.data = data
"""

# ==================================================================================================
class File:
    __slots__ = ('useCase', 'path', 'win32_path', 'dirname', 'name', 'extension', 'MIMEtype', \
                 'metadataFile', 'JSONmetadata', 'defaultDirAttributes', '_AreAttributesSet', '_occurredCriticalException')
    
    def __init__(self, path="", useCase="", defaultDirAttributes={}, JSONmetadata={}):
        self.useCase = JSONmetadata["Use-Case$Folder"]          if (JSONmetadata) \
            else useCase # use case of the file: it should be the "folder" where the file it's located.
        self.path    = JSONmetadata["General-Metadata"]["ABSOLUTEPATH"] if ( (JSONmetadata) and ("ABSOLUTEPATH" in JSONmetadata["General-Metadata"]) ) \
            else       JSONmetadata["General-Metadata"]["PATH"        ] if ( (JSONmetadata) and ("PATH"         in JSONmetadata["General-Metadata"]) ) \
            else path # full path of the file as a Unicode string ('utf-8'). # "ABSOLUTEPATH" is always expected in "General-Metadata" for main files/containers.
        self.win32_path = ""
        self.dirname = ""
        self.name      = JSONmetadata["General-Metadata"]["FILENAME"]  if (JSONmetadata) else ""
        self.extension = JSONmetadata["General-Metadata"]["EXTENSION"] if (JSONmetadata) else ""
        self.MIMEtype  = JSONmetadata["General-Metadata"]["TYPE"]      if (JSONmetadata) else ""

        '''
        *self.metadataFile* points to the metadata file generated by the "NLNZ Metadata Extractor v3.0" tool.
        The tool generates an XML document for each document/file, using the following process specs:
        - Configuration type: "Extract in Native form"
        - Destination: "../_metadata/<dataset>"
        - Profile: "Default"
        - Run as a simple (flatten structure) object.
        '''
        self.metadataFile = "" # location of the "general" metadata file.
        
        self.JSONmetadata = JSONmetadata
        self.defaultDirAttributes = defaultDirAttributes
        self._AreAttributesSet = False # a check flag to see if the attributes have been set.
        self._occurredCriticalException = False
        if (self.JSONmetadata):
            return # don't set the Dir. attributes.
        if os.path.isfile(path):
            self.setAttributes()
        else:
            Utils.output(f"'{path}' is not a file.", _error=True, _print=True)


    def setAttributes(self):
        if os.path.isfile(self.path):
            _, self.extension = os.path.splitext(self.path)
            self.dirname = os.path.dirname(self.path)
            self.name = os.path.basename(self.path)
            self.extension = "" + self.extension.translate({ord(c): None for c in '.'})
            self.MIMEtype = Utils._config["MIME-Types"].get(self.extension.upper(), mimetypes.MimeTypes().guess_type(self.name)[0])
            if ((self.MIMEtype is None) or not(self.MIMEtype)):
                Utils.output(f"MIME type unrecognisable for file extension: *{self.extension}*", _error=True, _print=True)
            self.metadataFile = (Utils.getFolder("Metadata-Folder") + self.name + Utils._config["Metadata-File-Extension"])
            if (platform.system() == 'Windows'):
                self.win32_path = self.path[4:] # The *win32* API does not support the WinAPI prefix for long paths. We ignore the WinAPI prefix.
                if (len(self.path) > Utils.WinAPI_MAX_PATH):
                    Utils.output(\
    f"""WARNING => Path length: {len(self.path)} >> {(len(self.path) / Utils.WinAPI_MAX_PATH * 100):.2f}% (max. path length: {Utils.WinAPI_MAX_PATH})""", _print=True)
                    # For .doc and .msg file formats, we use the *win32* API:
                    if ((self.extension.upper() == "DOC") or (self.extension.upper() == "MSG")):
                        # copy original file in a temporary location with a random integer + remove the WinAPI prefix:
                        self.win32_path = (Utils.getFolder("Temp-Folder")[4:] + "_MEL-" + str(random.randint(1,1000)) + "." + self.extension)
                        Utils.output(f"""The *win32* API does not support paths greater or equal to {
                            Utils.WinAPI_MAX_PATH}.\nPATH=[{self.path[4:]}]\nLENGTH={len(self.path[4:])}, ({
                            len(self.path[4:]) - Utils.WinAPI_MAX_PATH} longer)\n*** Copying original file to a shorter temporary path:\nTEMP=[{
                            self.win32_path}]\n""", _print=True)
                        shutil.copy2(self.path, self.win32_path) # we use a temporary copy; hopefully, it's a shorter path.
            else:
                self.win32_path = self.path # On non-windows platforms
            self._AreAttributesSet = True
            Utils.output(
                f"""Directory: [{self.dirname}]
                    File Name: [{self.name}]
                    Extension: [{self.extension}, {self.MIMEtype}]
                    Metadata File: [{self.metadataFile}]""", _print=False)
        else:
            self._AreAttributesSet = False
            Utils.output(f"'{self.path}' is not a file.", _error=True, _print=True)


    def print(self):
        Utils.output(f"""
            * Use Case: {self.useCase}
            * Path: {self.path}
            * Directory: {self.dirname}
            * File Name: {self.name}
            * Extension: {self.extension}
            * MIME Type: {self.MIMEtype}
            * Metadata File: {self.metadataFile}
            * Default Directory Attributes: {json.dumps(self.defaultDirAttributes, indent=4)}
            * JSON Metadata: {json.dumps(self.JSONmetadata, indent=4)}
            * Are attributes set?: {self._AreAttributesSet}
            """)


    def isEmpty(self):
        if (self.defaultDirAttributes.get("FILELENGTH", -1) == 0): # If the attribute is not found, then we use a default value of -1.
            Utils.output(f"The file size is 0K: FILE EMPTY NOT PROCESSABLE!", _print=True)
            return True
        return False


    def isProtected(self):
        if ((self.defaultDirAttributes.get("FILELENGTH", -1) <= 256) and (self.name[0:2] == "~$")): # a system/hidden/protected file
            Utils.output(f"The file is protected: FILE NOT PROCESSABLE!", _print=True)
            return True
        return False


    def occurredCriticalException(self):
        return self._occurredCriticalException


    def extractMetadata(self, _perform_NER_processing=None):
        if (self.JSONmetadata): # if we already have the metadata structure: for NLP/NER processing cases or metadata prunning on CouchDB...
            return self.JSONmetadata
        if (self.isEmpty() or self.isProtected()):
            return {}
        if (self.useCase == ""):
            Utils.output(f"The use case label is empty.", _error=True, _print=True)
            return {}
        absolute_path = ""
        read_ok = False
        try:
            tree = ET.parse(os.path.abspath(self.metadataFile))
            read_ok = True
        except: # if ET.parse fails: ignore and return an empty dictionary.
            read_ok = False
        if (read_ok):
            root = tree.getroot()
            for item in root.findall('./*'): # xml to json simple conversion
                self.JSONmetadata["General-Metadata" if (item.tag == "METADATA") else item.tag] = \
                    parker.data(item.findall('./*')) if item else "" # checks if element is empty
            ''' Checks if the "ABSOLUTEPATH" field from the general metadata info. corresponds to the actual *self.path*.
            This is done for the cases were we have filenames with the same name but under different folders.
            We take into account if there's a difference in the paths due to *extended Windows paths* (Utils.WinAPI_ExtendedLengthPathPrefix). '''
            if (self.JSONmetadata["General-Metadata"]["ABSOLUTEPATH"][:3] == Utils.WinAPI_ExtendedLengthPathPrefix):
                absolute_path = self.JSONmetadata["General-Metadata"]["ABSOLUTEPATH"] \
                    if  (self.path[:3] == Utils.WinAPI_ExtendedLengthPathPrefix) \
                    else self.JSONmetadata["General-Metadata"]["ABSOLUTEPATH"][4:]
            else:
                absolute_path = self.JSONmetadata["General-Metadata"]["ABSOLUTEPATH"]
            # if (self.path != absolute_path) ==> assign default general metadata values.
            # this indicates that the general metadata file does not correspond to the actual file (but they have the same filename).
            read_ok = (self.path == absolute_path)
        if not (read_ok): # By default...
            self.JSONmetadata = {} # clears full JSON object.
            self.JSONmetadata["General-Metadata"] = {} # resets general-metadata section.
            by_default = self.defaultDirAttributes # set the default directory attributes
            if not by_default: # if there are not any directory attributes set by default: special case for the content of zip and msg files.
                by_default["PATH"] = self.path
                by_default["FILENAME"] = self.name
                by_default["EXTENSION"] = self.extension
                by_default["TYPE"] = self.MIMEtype
            self.JSONmetadata["General-Metadata"] = by_default
        # Adding general attributes and correcting values for the Metadata Extraction Tool:
        self.JSONmetadata["Use-Case$Folder"] = self.useCase
        self.JSONmetadata["General-Metadata"]["TYPE"] = self.MIMEtype
        self.JSONmetadata["Specific-Metadata"] = self.extract_ext_info()
        self.JSONmetadata["Associated-Metadata"] = self.extract_associated_info()
        # perform NLP-NER processing? (1st condition) Default action; (2nd condition) Force processing
        if (( (_perform_NER_processing is None) and (Dataset.NER_Complete()) ) or \
              (_perform_NER_processing == True)):
            self.JSONmetadata.update( NER.process(self.JSONmetadata) )
        Utils.output(json.dumps(self.JSONmetadata, indent=4), _print=False)
        return self.JSONmetadata


    def extractPDFinfo(self):
        pdfMetadata, pdfOutline, props_exception_details, txt_exception_details = {}, {}, {}, {}

        def textExtractionException(e_type, e_details):
            txt_exception_details["Type"] = e_type
            txt_exception_details["Details"] = e_details
            Utils.output(f"An Exception occurred while processing the PDF file *text content*: " + \
                         json.dumps(txt_exception_details, indent=4), _print=True)
            pdfMetadata["PDF.Text-Processing-Exception"] = txt_exception_details

        def translateChar(c):
            return "_" if (ord(c) > 122) else c # if (ord(c) > ord('z'))
        
        def getStrInEnglishCharacterSet(_str):
            trans = ""
            for c in _str:
                trans = trans + translateChar(c)
            return trans

        def rpl_dqWithS(matchObj): # \1\g<sub>\5\6
            return (matchObj.group(1) + \
                    matchObj.group("sub").replace('"s', "'s").replace('"S', "'S") +\
                    matchObj.group(5) + matchObj.group(6))
        
        def rpl_sWithDQ(matchObj): # \1\g<sub>\7\8
            return (matchObj.group(1) + \
                    matchObj.group("sub").replace('s"', "s'").replace('S"', "S'") +\
                    matchObj.group(7) + matchObj.group(8))
        
        def rpl_dqList(matchObj): # \1\2\g<sub>\10\11
            return (matchObj.group(1) + matchObj.group(2) +\
                    matchObj.group("sub").replace('"', "'") +\
                    matchObj.group(10) + matchObj.group(11))
        
        properties_read_ok, text_read_ok  = False, False
        strOutline, text = "", ""
        numberOfPages = 0
        try: # Extract PDF document properties:
            with open(self.path, 'rb') as f:
                pdf = PdfFileReader(f)
                information = pdf.getDocumentInfo()
                numberOfPages = pdf.getNumPages()
                strOutline = pdf.outlines.__str__()
                properties_read_ok = True
        except Exception as e: # if any failure occurs: return a dictionary with the Exception details.
            props_exception_details["Type"] = str(type(e))
            props_exception_details["Details"] = str(e)
            Utils.output(f"An Exception occurred while processing the PDF file *properties*: " + \
                         json.dumps(props_exception_details, indent=4), _print=True)
            pdfMetadata["PDF.Properties-Processing-Exception"] = props_exception_details
            if (str(e) == "file has not been decrypted"):
                self._occurredCriticalException = True
                Directory.encrypted_PDFs.append(self.path)
            properties_read_ok = False
        
        ''' # Extract PDF text:
        currentPageNumber = 0
        while (currentPageNumber < numberOfPages):
            pdfPage = pdf.getPage(currentPageNumber)
            currentPageNumber += 1
            text = text + pdfPage.extractText()
        * pdfPage.extractText() is a faulty method!
        * References:   <https://stackoverflow.com/questions/4203414/pypdf-unable-to-extract-text-from-some-pages-in-my-pdf>
        * Solution:     <https://en.wikipedia.org/wiki/Pdftotext>
        * Switching to: <xpdf-tools-win-4.02\bin64\pdftotext>
        * Option: "-layout" -> Maintain (as best as possible) the original physical layout of the text.
        * (1st) try: simple text extraction without preserving the document layout --> content-oriented. (... if it fails then try with ...)
        * (2nd) try: text extraction preserving the layout: "-layout" argument     --> layout-oriented.
        * Problem on handling UNICODE (non-English) characters in filenames:
        pdftotext doesn't support UNICODE characters in filenames.  Implemented a workaround.
        '''
        content_oriented_extraction = ["pdftotext",            self.path, (Utils.getFolder("Temp-Folder") + getStrInEnglishCharacterSet((self.name +        '.txt')))]
        layout_oriented_extraction  = ["pdftotext", "-layout", self.path, (Utils.getFolder("Temp-Folder") + getStrInEnglishCharacterSet((self.name + "-layout.txt")))]
        arg_layout = "-layout" if (Dataset.pdftotext_preserveLayout()) else "" # 1st try: based on the set flag
        try_again = True
        for attempt in range(2): # two attempts: [0, 1]
            try:
                command = []
                if (try_again): # control flow of attempts
                    command = content_oriented_extraction if not(arg_layout) else layout_oriented_extraction
                    temp_txt_filename = command[len(command)-1]
                    Utils.output(f"""pdftotext command (attempt #{attempt+1}):\n{command}""", _print=False)
                    subprocess.run(command, encoding="utf-8")
                    with open(temp_txt_filename, encoding="utf-8") as f_txt: # mode="rt" (default)
                        text = f_txt.read().encode("utf-8").decode("utf-8")
            except Exception as e: # if any failure occurs: return a dictionary with the Exception details.
                if ((attempt == 0) and ("'utf-8' codec can't decode byte 0x" in str(e))): # not a suitable result with "-layout".
                    arg_layout = "" if (Dataset.pdftotext_preserveLayout()) else "-layout" # 2nd try: flip values
                    try_again = True
                else: # any other exception
                    textExtractionException(str(type(e)), str(e))
                    try_again = False
                    text_read_ok = False
            else:
                text_read_ok = True
                break # exit the loop with success
        else: # all attempts fail
            text_read_ok = False
        # Utils.output(f"PDF text extraction:\n{text}\n", _print=True)
        try:
            if (text_read_ok):
                EMPTY = ""
                if (numberOfPages > 0):
                    for _ in range(numberOfPages):
                        EMPTY += chr(12) # "Form Feed (FF)" character
                    Utils.output(f"... Num-pages={numberOfPages} | EMPTY={EMPTY}", _print=False)
                if (text == EMPTY): # If can not extract text then use OCR lib to extract the scanned PDF file.
                    text = textract.process(self.path, method='tesseract', encoding='utf-8')
        except Exception as e: # if any failure occurs: return a dictionary with the Exception details.
            textExtractionException(str(type(e)), str(e))
            text_read_ok = False
        if (properties_read_ok):
            try: # *information.author* might not exist:
                pdfMetadata["PDF.Author"  ] = (information.author   if (information.author is not None) else "")
                pdfMetadata["PDF.Creator" ] = (information.creator  if (information.author is not None) else "")
                pdfMetadata["PDF.Producer"] = (information.producer if (information.author is not None) else "")
                pdfMetadata["PDF.Subject" ] = (information.subject  if (information.author is not None) else "")
                pdfMetadata["PDF.Title"   ] = (information.title    if (information.author is not None) else "")
            except:
                '''
When reading the attribute ".author", it seems that sometimes, it might not even exist in the "information" object:
## pdfMetadata["PDF.Author"] = (information.author if not(information.author is None) else "")
## AttributeError: 'NoneType' object has no attribute 'author'
We ignore this type of failure...'''
                pdfMetadata["PDF.Author"  ] = ""
                pdfMetadata["PDF.Creator" ] = ""
                pdfMetadata["PDF.Producer"] = ""
                pdfMetadata["PDF.Subject" ] = ""
                pdfMetadata["PDF.Title"   ] = ""
            pdfMetadata["PDF.Number-of-pages"] = numberOfPages
            '''
The following replacements are needed to convert the PDF Outline into a correct JSON object (format/syntax).
They have been tested and implemented on the go (as needed) based on the different cases.
The order of replacements IS IMPORTANT and may incurred in failures if they are changed.'''
            Utils.output(f"Outline's JSON Structure -- original:\n{strOutline}", _print=False)
            strOutline = strOutline\
                .replace("'", '"')\
                .replace(' "$" ', " '$' ")\
                .replace(' "s ', "'s ")\
                .replace(u'\u2002', " ")\
                .replace('\\u2002"', '')\
                .replace('""', '"')\
                .replace("/", '')\
                .replace("\\", '\\\\')\
                .replace('\\\\"s', "'s")\
                .replace(': b"', ': "-(binary)-')\
                .replace("IndirectObject(", '"IndirectObject(')\
                .replace(", 0)", ', 0)"')\
                .replace("<PyPDF4", '"<PyPDF4')\
                .replace(">", '>"')\
                .replace('>"."', '>."')\
                .replace('>". ', '>. ')\
                .replace('>"",', '>",')\
                .replace('; >"', "; >")\
                .replace(u"\u03b1", "-'alpha'")\
                .replace(u"\u03b2", "-'beta'")\
                .replace(u"\u03b3", "-'gamma'")\
                .replace(u"\u03b4", "-'delta'")\
                .replace(u"\u03bc", "-'mu'")\
                .replace(u"\u2714", "-(check-mark)")\
                .replace('"Title": ", "Page":', '"Title": "", "Page":')
            '''
            Testing reference: <https://www.regexpal.com/>, <https://regex101.com/>
            '''
            strOutline = re.sub(r'( |")([A-Z])"(\w+) ', r" \1\2'\3 ", strOutline) # D"Ambrosio , O"Reilly ...
            strOutline = re.sub(r'(\w+)(, "Page": ")', r'\1"\2', strOutline)
            strOutline = re.sub(r'"(\w+\.\.\.)(", "Page":)', r"'\1\2", strOutline)
            strOutline = re.sub(r'("Title": ")(?P<sub>([^"]+"(s|S))+)([^"]+)(", "Page":)', rpl_dqWithS, strOutline)
            strOutline = re.sub(r'(\w+)(\s+)>"(\s+)?(.)?(.)?(\d+)', r"\1\2>\3\4\5\6", strOutline)
            strOutline = re.sub(r'(\s+)\(>"(\s+)?(\d+)(\s+)?\)', r"\1(>\2\3\4)", strOutline)
            strOutline = re.sub(r'(\w+) "(\w+)"(-)?(\w+)(.)? ', r"\1 '\2'\3\4\5 ", strOutline)
            strOutline = re.sub(r'("Title": ")(?P<sub>([^"]+(s|S)"(\)|,|\.)? (\w+))+)([^"]+)(", "Page":)', rpl_sWithDQ, strOutline)
            strOutline = re.sub(r'(\w+) "((\w|\s)+)"(.?)(\w+)?', r"\1 '\2'\4\5", strOutline)
            strOutline = re.sub(r' "(\w+ \- (\w|\s)+)"( |[:,])', r" '\1'\3", strOutline)
            strOutline = re.sub(r'("Title": ")(\(\w\) )?([^"]+)"((\w|\s)+s\')', r"\1\2\3'\5", strOutline)
            strOutline = re.sub(r'("Title": ")([^"]+)(?P<sub>((")(([^"]+))(")(, )?)+)([^"]+)(", "Page":)', rpl_dqList, strOutline)
            strOutline = re.sub(r'(-)(")(\w+)', r"\1'\3", strOutline) # Meta-"omics
            # Special cases:
            # '; >"', "; >": NHRMC Grants
            # '(\w+)"s ' , "(\w+)'s ": NHRMC Grants
            # '(\w+) >" ', "(\w+) > ": NHRMC Grants
            # '"(\w+)"-', "'(\w+)'-": NHRMC Grants
            # '>"."', '>."': DoEE-Species
            # '>". "', '>. "': DoEE-Species
            # '(>"1000)', '(>1000)': DoEE-Species
            # '(\w+)s" (\w+)', "(\w+)s' (\w+)": DoEE-Species
            # '(\w+) "((\w|\s)+)" (\w+)', "(\w+) '((\w|\s)+)' (\w+)": DoEE-Species
            # \u03b1: GREEK SMALL LETTER ALPHA
            # \u03b2: GREEK SMALL LETTER BETA
            # \u03b3: GREEK SMALL LETTER GAMMA
            # \u03b4: GREEK SMALL LETTER DELTA
            # \u03bc: GREEK SMALL LETTER MU
            # \u2714: HEAVY CHECK MARK
            # \u2002: EN SPACE
            Utils.output(f"Outline's JSON Structure -- replacements:\n{strOutline}", _print=False)
            pdfOutline = json.loads(strOutline) if (Utils.isJSON(strOutline)) else strOutline.replace('"', "'")
            pdfMetadata["PDF.Outlines"] = pdfOutline
        if (text_read_ok):
            _txt = Text(text)
            pdfMetadata["text-analysis"] = _txt.analysis()
        Utils.output(json.dumps(pdfMetadata, indent=4), _print=False)
        return pdfMetadata


    # Applies for DOCX and PPTX files.
    def extractMSOfficeFileProps(self, office_file):
        cp = office_file.core_properties
        return {
            "author": cp.author,
            "category": cp.category,
            "comments": cp.comments,
            "content_status": cp.content_status,
            "created": f"{Utils.getFormattedDT(cp.created)}",
            "identifier": cp.identifier,
            "keywords": cp.keywords,
            "language": cp.language,
            "last_modified_by": cp.last_modified_by,
            "last_printed": (f"{Utils.getFormattedDT(cp.last_printed)}" if (cp.last_printed is not None) else ""),
            "modified": f"{Utils.getFormattedDT(cp.modified)}",
            "revision": cp.revision,
            "subject": cp.subject,
            "title": cp.title,
            "version": cp.version
        }


    def extractPPTXinfo(self):
        # For all the content in the presentation:
        bolds, italics, underlines, hyperlinks = {}, {}, {}, {}
        # "text" will be populated with a list of strings, one for each text placeholder and shape (text_frame), for each slide of the presentation.
        text = []
        ''' =============================================================================================================================== '''        
        # Paragraphs from TextFrame objects:
        def extractParagraphs(paragraphs, check_in_text, delimiter=" "):
            text_paras = ""
            if len(paragraphs) > 0:
                for para in paragraphs:
                    for run in para.runs: # to get specific styles:
                        if run.font.bold : # bold style
                            Utils.addItemInCountingList(bolds, run.text)
                        if run.font.italic : # italic style
                            Utils.addItemInCountingList(italics, run.text)
                        if run.font.underline : # underline style
                            Utils.addItemInCountingList(underlines, run.text)
                        if run.hyperlink :
                            Utils.addItemInCountingList(hyperlinks, run.hyperlink.address)
                        if not (check_in_text): # appends the additional text without checking.
                            text.append(run.text)
                        else: # checks if the text already exits...
                            if not (run.text in " ".join(t for t in text)): # not found in the overall text of the current slide.
                                text.append(run.text)
                        text_paras += (delimiter if (len(text_paras) > 0) else "") + run.text
            return text_paras
        ''' =============================================================================================================================== '''
        def extractShapes(shapes):
            imagesMetadata = []
            for s in shapes:
                if (s.shape_type == MSO_SHAPE_TYPE.PICTURE):
                    image = s.image
                    imagesMetadata.append({
                        "content_type": image.content_type,
                        "dpi": image.dpi,
                        "extension": image.ext,
                        "filename": image.filename,
                        "size": image.size
                    })
                # text in paragraphs:
                if not s.has_text_frame:
                    continue
                extractParagraphs(s.text_frame.paragraphs, True)
            return imagesMetadata
            
        ''' =============================================================================================================================== '''
        prs = Presentation(self.path) # PPTX file: the presentation
        prs_props = self.extractMSOfficeFileProps(prs)

        slides, num_slides = [], 0
        for slide in prs.slides:
            num_slides += 1
            # initialize data attributes/properties for the new slide to process:
            slideMetadata, subtitle, text, notes_slide, num_shapes_in_notes = {}, "", [], "", 0
            tablesMetadata, dict_tables, text_tables, t_count = {}, {}, "", 0 # for table processing...
            chartsMetadata, dict_charts, text_charts, chart_count = {}, {}, "", 0 # for chart processing...
            imagesMetadata = [] # for image processing...
            for placeholder in slide.placeholders:
                if placeholder.has_text_frame:
                    text.append(placeholder.text_frame.text)
                if ("SUBTITLE" in str(placeholder.placeholder_format.type).upper()):
                    subtitle = placeholder.text_frame.text
                    
                ''' =============================================================================================================================== '''
                ''' # /*- THE IMPLEMENTATION IS NOT COMPLETE /*-
                    # The current implementation is missing the extraction of some supported components from the OBJECT MODEL, such as:
                    # images and tables inserted as OLE OBJECTS (from Excel or Word documents).
                    # REFERENCE: <https://python-pptx.readthedocs.io/en/latest/api/chart.html> '''
                if (placeholder.has_chart): # text in charts: not implemented!
                    chart_count += 1
                    chart = placeholder.chart
                    chartMetadata, chartTitle, chartCategoryAxis = {}, "", ""
                    if (chart.has_title):
                        chartTitle = extractParagraphs(chart.chart_title.text_frame.paragraphs, False) \
                                    if (chart.chart_title.has_text_frame) else "" 
                    if (chart.category_axis):
                        axis = chart.category_axis
                        chartCategoryAxis = extractParagraphs(axis.axis_title.text_frame.paragraphs, False) \
                                    if (axis.has_title) else "" 
                    chartMetadata = {
                        "title": chartTitle,
                        "category-axis": chartCategoryAxis
                    }
                    text_charts += f"[chart={chart_count}]:" + \
                                    "\n" + chartMetadata["title"] + \
                                    "\n" + chartMetadata["category-axis"] + "\n\n"
                    dict_charts[f"chart-{chart_count}"] = chartMetadata
                    chartsMetadata = {
                        'structure': dict_charts,
                        'flatten': text_charts
                    }
                    Utils.output(f"<!!> *Processing Charts* HAS NOT BEEN IMPLEMENTED COMPLETELY <!!>", _print=True)
                    Utils.output(json.dumps(chartsMetadata, indent=4), _print=True)

                ''' =============================================================================================================================== '''
                if (placeholder.has_table): # text in tables:
                    t = placeholder.table
                    t_count += 1
                    rows, r_count = [], 0 # initialize data attributes/properties for the new table to process.
                    for r in t.rows:
                        r_count += 1
                        cells, c_count = [], 0 # initialize data attributes/properties for the new row to process.
                        for c in r.cells:
                            c_count += 1
                            c_text = extractParagraphs(c.text_frame.paragraphs, False) # text extraction from the paragraphs collection.
                            cells.append(c_text)
                            text_tables += f"[t={t_count},r={r_count},c={c_count}]:" + c_text + "\n"
                        text_tables += "\n"
                        rows.append(cells)
                    dict_tables[f"table-{t_count}"] = rows
                    tablesMetadata = {
                        'structure': dict_tables,
                        'flatten': text_tables
                    }
            
            ''' =============================================================================================================================== '''
            imagesMetadata = extractShapes(slide.shapes)
            if (slide.has_notes_slide):
                notes_slide = extractParagraphs(slide.notes_slide.notes_text_frame.paragraphs, False, delimiter="\n")
                extractShapes(slide.notes_slide.shapes)
                num_shapes_in_notes = len(slide.notes_slide.shapes)
            slideMetadata = {
                "slide-number": num_slides,
                "title": slide.shapes.title.text if (slide.shapes.title is not None) else "",
                "sub-title": subtitle,
                "number-of-shapes": len(slide.shapes),
                "tables": tablesMetadata,
                "images": imagesMetadata,
                "text": text,
                "notes": notes_slide,
                "number-of-shapes-in-notes": num_shapes_in_notes
            }
            slides.append(slideMetadata)
        pptxMetadata = {
            "document-properties": prs_props,
            "number-of-slides": num_slides,
            "slides": slides,
            'bold-phrases': bolds,
            'italic-phrases': italics,
            'underline-phrases': underlines,
            'hyperlinks': hyperlinks
        }
        # Extracts all the content (text) for its analysis:
        _txt = Text(" \n\n\n ".join(t for t in 
                                    map(lambda s: " \n ".join(t for t in s["text"]), pptxMetadata["slides"])
                    ))
        pptxMetadata["text-analysis"] = _txt.analysis()
        Utils.output(json.dumps(pptxMetadata, indent=4), _print=False)
        return pptxMetadata


    def extractDOCXinfo(self):
        # For all the content:
        headings, bolds, italics, underlines = {}, {}, {}, {}

        def extractParagraphs(paragraphs):
            doc_text_paras = ""
            if len(paragraphs) > 0:
                for para in paragraphs:
                    # Find email and phone numbers within the paragraph text:
                    text = para.text
                    doc_text_paras += ("\n\n" if (len(doc_text_paras) > 0) else "") + para.text
                    if (para.style.name != "Normal"):
                        if (headings.get(para.style.name, "") == ""):
                            headings[para.style.name] = [] # creates the list
                        headings[para.style.name].append(text)
                    for run in para.runs:
                        if run.bold : # bold style
                            Utils.addItemInCountingList(bolds, run.text)
                        if run.italic : # italic style
                            Utils.addItemInCountingList(italics, run.text)
                        if run.underline : # underline style
                            Utils.addItemInCountingList(underlines, run.text)
            return doc_text_paras

        def extractTables(tables):
            tablesMetadata = {}
            if len(tables) > 0:
                dict_tables, rows, cells = {}, [], []
                doc_text_tables = "\n"
                t_count, r_count, c_count = 0, 0, 0
                for t in tables:
                    t_count += 1
                    r_count, rows = 0, []
                    for r in t.rows:
                        r_count += 1
                        c_count, cells, cell = 0, [], {} # each cell may be composed of paragraphs and tables.
                        for c in r.cells:
                            c_count += 1
                            c_text = extractParagraphs(c.paragraphs) # text extraction from the paragraphs collection.
                            cell["text"] = c_text
                            cell["tables"] = extractTables(c.tables)
                            cells.append(cell)
                            doc_text_tables += f"[t={t_count},r={r_count},c={c_count}]:" + c_text + "\n"
                        doc_text_tables += "\n"
                        rows.append(cells)
                    dict_tables[f"table-{t_count}"] = rows
                tablesMetadata = {
                    'structure': dict_tables,
                    'flatten': doc_text_tables
                }
            return tablesMetadata

        def extractSectionPart(part): # header | footer
            return {
                'text-paragraphs': extractParagraphs(part.paragraphs),
                'tables': extractTables(part.tables),
                'all-text': extractParagraphs(part.paragraphs) +\
                    "" if (len(part.tables) == 0) else extractTables(part.tables)["flatten"]
            }

        document = Document(self.path)
        doc_props = self.extractMSOfficeFileProps(document)
        rels = document.part.rels
        hyperlinks, list_of_hyperlinks = [], ""
        for rel in rels:
            if rels[rel].reltype == RT.HYPERLINK:
                list_of_hyperlinks += rels[rel]._target
                hyperlinks.append(rels[rel]._target)
        sections, txt_sections, s = {}, "", 0
        for section in document.sections:
            s += 1
            header = extractSectionPart(section.header)
            footer = extractSectionPart(section.footer)
            txt_sections += header["all-text"] + " -||- " + footer["all-text"]
            sections[f"section-{s}"] = {
                'header': header,
                'footer': footer
            }
        docxMetadata = {
            'document-properties': doc_props,
            'headings': headings,
            'text-paragraphs': extractParagraphs(document.paragraphs),
            'bold-phrases': bolds,
            'italic-phrases': italics,
            'underline-phrases': underlines,
            'tables': extractTables(document.tables),
            'hyperlinks': hyperlinks,
            'sections': sections
        }
        _txt_flatten_tables = "" if (len(document.tables) == 0) else docxMetadata["tables"]["flatten"]
        _txt = Text(docxMetadata["text-paragraphs"] + " -||- " +\
                    _txt_flatten_tables + " -||- " +\
                    txt_sections + " -||- " +\
                    list_of_hyperlinks)
        docxMetadata["text-analysis"] = _txt.analysis()
        Utils.output(json.dumps(docxMetadata, indent=4), _print=False)
        return docxMetadata


    def extractDOCMinfo(self):
        o = os.popen(f'{Utils._config["ConvertDOCMtoDOCX-.NET-Utility"]} "' + self.path + '"').read() # runs the utility to convert .docm to .docx
        Utils.output(o, _print=False)
        for line in o.split('\n'):
            match = re.search(r'\* OUTPUT: ', line) # gets the output file
            if (match):
                outputfile = line.split(': ')
        Utils.output(outputfile, _print=False)
        doc_original_path, doc_temp_path = self.path, outputfile[1]
        self.path = doc_temp_path # sets the document path to the temporal .docx created by the utility.
        docmMetadata = self.extractDOCXinfo()
        self.path = doc_original_path # restores original path. 
        return docmMetadata


    def extractOLEinfo(self):
        oleMetadata, omOLEfile, omParsingIssues = {}, {}, {}
        if (olefile.isOleFile(self.path)): # only files that had been detected as "OLE".
            o = os.popen('olemeta "' + self.path + '"').read() # runs the olemeta utility
            _property = ""
            for line in o.split('\n'):
                match = re.search(r'\|[\w.-]+', line)
                if (match):
                    key_value = line.split('|') # # each line has the structure: "|Property             |Value                         |"
                    if (key_value[1].strip() != "Property"):
                        if (key_value[1].strip() != ""):
                            _property = key_value[1].strip()
                            oleMetadata[_property] = key_value[2].strip()
                        else: # append the additional line as part of the value.
                            oleMetadata[_property] = oleMetadata[_property] + " " + key_value[2].strip()
            text_content = ""
            with olefile.OleFileIO(self.path) as ole: # perform all operations on the ole object
                if ole.parsing_issues:
                    for exctype, msg in ole.parsing_issues:
                        omParsingIssues[exctype.__name__] = msg
                omOLEfile["Parsing-Issues"] = omParsingIssues
                streams = ole.listdir() # all the content | (streams=True)
                omOLEfile["Streams-and-Storages-List"] = streams.copy() # copy the list
                # For specific legacy file formats (DOC or XLS), we extract __only__ the main content stream.
                streams = ['WordDocument'] if (self.extension.upper() == "DOC") \
                    else  ['Workbook'    ] if (self.extension.upper() == "XLS") \
                    else streams
                for stream in streams:
                    stream_buffer = ole.openstream(stream).read()
                    text_content += f'<stream name="{stream}">' + \
                                    "".join(map(chr, stream_buffer)) + f'</stream>' # from an array of bytes to string
            oleMetadata["OLEfile-Structure"] = omOLEfile
            if (text_content):
                txt = Text(text_content)
                oleMetadata["text-analysis"] = txt.analysis()
        else: # not an OLE file.
            '''
            Basic processing for *.DOC files
            '''
            if (self.extension.upper() == "DOC"):
                if (platform.system() != 'Windows'):
                    Utils.output(f"The method uses a *win32* system library.\nCannot run on a {platform.system()} platform.", _error=True, _print=True)
                    return
                word = win32.gencache.EnsureDispatch("Word.Application")
                word.Visible = False
                # we used the win32_path attribute: points to the corrected path (might be a temp. location)
                word.Documents.Open(self.win32_path if (platform.system() == 'Windows') else self.path)
                word.ActiveDocument.ActiveWindow.View.Type = 3  # prevent that Word opens itself
                doc = word.ActiveDocument
                txt = Text(doc.Range().Text)
                oleMetadata["text-analysis"] = txt.analysis()
                word.Quit()
                del word
            else:
                details = \
    f"""The file has an extension of "{self.extension}" so it should be an OLE file but it is not.
    MEL was not able to process the file. The file may be corrupt or of an incorrect format."""
                exception = {}
                exception["Details"] = details
                Utils.output(f"An Exception occurred while processing the {self.extension} file: " + \
                             json.dumps(exception, indent=4), _print=True)
                oleMetadata["OLE.Processing-Exception"] = exception
        Utils.output(json.dumps(oleMetadata, indent=4), _print=False)
        return oleMetadata


    def extractMSGinfo(self):
        if (platform.system() != 'Windows'):
            Utils.output(f"The method uses a *win32* system library.\nCannot run on a {platform.system()} platform.", _error=True, _print=True)
            return
        msgMetadata, attachments = {}, {}
        outlook = win32.Dispatch("Outlook.Application").GetNamespace("MAPI")
        ''' If we have the following exception, it could be that the file size is 0K (empty file).

## msg = outlook.OpenSharedItem(self.path)
## File "<COMObject GetNamespace>", line 2, in OpenSharedItem
## pywintypes.com_error: (-2147352567, 'Exception occurred.', 
(4096, 'Microsoft Outlook',
"We can't open 'E:\\_temp\\DepFin-Project\\DoEE_endangered-species\\SpeciesDocumentSets\\Mammals - Gy...'. 
It's possible the file is already open, or you don't have permission to open it.\n\n
To check your permissions, right-click the file folder, then click Properties.", None, 0, -2147286960), None)
        
        In order to monitor better the metadata extraction process, we won't make any exception handling nor validation.
        '''
        if (self.isEmpty()):
            pass
        # we used the win32_path attribute: points to the corrected path (temp. location)
        msg = outlook.OpenSharedItem(self.win32_path if (platform.system() == 'Windows') else self.path)
        ''' /*- Problem regarding processing *.OFT files -- {ongoing} /*- '''

        msgMetadata["sender-name"] = msg.SenderName
        msgMetadata["sender-email-address"] = msg.SenderEmailAddress
        msgMetadata["sent-on"] = Utils.getFormattedDT(msg.SentOn)
        
        msgMetadata["creation-time"] = Utils.getFormattedDT(msg.CreationTime)
        msgMetadata["last-modification-time"] = Utils.getFormattedDT(msg.LastModificationTime)
        msgMetadata["importance"] = msg.Importance

        '''
When reading some of the following attributes, the library might fail with the following description:
## File "E:\Dropbox\Library\Software\Python\Python37\lib\site-packages\win32com\client\dynamic.py", line 527, in __getattr__
## raise AttributeError("%s.%s" % (self._username_, attr))
## AttributeError: OpenSharedItem._________
We ignore this type of failure...'''
        try:
            msgMetadata["is-marked-as-task"] = msg.IsMarkedAsTask
        except:
            msgMetadata["is-marked-as-task"] = False
        try:
            msgMetadata["received-by-name"] = msg.ReceivedByName # the display name of the true recipient for the mail message.
        except:
            msgMetadata["received-by-name"] = "<Unable to retrieved it>"
        try:
            msgMetadata["received-on-behalf-of-name"] = msg.ReceivedOnBehalfOfName # the display name of the user delegated to represent the recipient for the mail message.
        except:
            msgMetadata["received-on-behalf-of-name"] = "<Unable to retrieved it>"
        try:
            msgMetadata["HTML-body"] = msg.HTMLBody
        except:
            msgMetadata["HTML-body"] = "<Unable to retrieved it>"
        
        msgMetadata["received-time"] = Utils.getFormattedDT(msg.ReceivedTime)
        msgMetadata["conversation-topic"] = msg.ConversationTopic
        msgMetadata["subject"] = msg.Subject
        msgMetadata["body"] = msg.Body
        _txt = Text(msg.Body)
        msgMetadata["text-analysis"] = _txt.analysis()
        
        # To/CC/BCC are extracted through the Recipients collection:
        #msgMetadata["to-name"] = msg.To
        #msgMetadata["CC-name"] = msg.CC
        #msgMetadata["BCC-name"] = msg.BCC
        rTo, rCC, rBCC = {}, {}, {}
        nTo, nCC, nBCC = 0, 0, 0
        for r in msg.Recipients:
            rProp = {}
            rProp["name"] = r.Name
            rProp["address"] = r.Address
            if (r.Type == 1): # To
                nTo += 1
                rTo[f"To-recipient-{nTo}"] = rProp
            elif (r.Type == 2): # CC
                nCC += 1
                rCC[f"CC-recipient-{nCC}"] = rProp
            else: # BCC
                nBCC += 1
                rBCC[f"BCC-recipient-{nBCC}"] = rProp
            #Utils._log.info(f"Type={r.Type}, Name={r.Name}, Address={r.Address}")
        msgMetadata["To"] = rTo
        msgMetadata["CC"] = rCC
        msgMetadata["BCC"] = rBCC
        
        count_attachments = msg.Attachments.Count
        if count_attachments > 0:
            for item in range(count_attachments):
                attachment = {}
                attachment["position"] = msg.Attachments.Item(item + 1).Position
                attachment["display-name"] = msg.Attachments.Item(item + 1).DisplayName
                '''
When reading the attribute ".FileName" or ".PathName", the library might fail with the following description:
## File "E:\Dropbox\Library\Software\Python\Python37\lib\site-packages\win32com\client\dynamic.py", line 516, in __getattr__
## ret = self._oleobj_.Invoke(retEntry.dispid,0,invoke_type,1)
## pywintypes.com_error: (-2147352567, 'Exception occurred.', (4096, 'Microsoft Outlook', 'Outlook cannot perform this action on this type of attachment.', None, 0, -2147467259), None)
We ignore this type of failure... and, unfortunately, we won't be able to retrieve the attachment...'''
                try:
                    attachment["filename"] = msg.Attachments.Item(item + 1).FileName
                except:
                    attachment["filename"] = ""
                try:
                    attachment["pathname"] = msg.Attachments.Item(item + 1).PathName
                except:
                    attachment["pathname"] = ""
                attachment["size"] = msg.Attachments.Item(item + 1).Size
                attachment["type"] = {
# https://docs.microsoft.com/en-us/dotnet/api/microsoft.office.interop.outlook.olattachmenttype?view=outlook-pia
# (yes... hard-coded!)
                    1: "ByValue", 
                    5: "EmbeddedItem",
                    6: "OLE"
                }.get(msg.Attachments.Item(item + 1).Type, 1)
                if (len(attachment["filename"]) > 0): # extracts the attachment file:
                    attachmentFileName = (Utils.getFolder("Temp-Folder") + attachment["filename"])
                    msg.Attachments.Item(item + 1).SaveAsFile(attachmentFileName)
                    attachmentFile = File(path=attachmentFileName, useCase=self.useCase, defaultDirAttributes={}, JSONmetadata={}) # default Directory attributes are set to empty.
                    attachment["metadata"] = attachmentFile.extractMetadata()
                else: # the filename is empty ==> cannot process the attachment.
                    attachment["metadata"] = {}
                attachments[f"attachment-{item + 1}"] = attachment
        msgMetadata["attachments"] = attachments
        del outlook, msg
        
        msgMetadata["OLEfile-Metadata"] = self.extractOLEinfo()
        Utils.output(json.dumps(msgMetadata, indent=4), _print=False)
        return msgMetadata


    def extractCSVinfo(self):
        """
        try:
            # ---
        except UnicodeEncodeError as err:
            Utils._log.info(err)
            #Utils.output(Utils.getValidEncoding(err), _print=True)
        v encoding='utf-8'
        v encoding='cp437' # Codepage 437 is the original DOS encoding.  All codes are defined.
        X encoding='cp1252' # Windows codepage
        X encoding="ISO-8859-1"
        X encoding='latin-1'
        X encoding='cp850'
        """
        csvRows = []
        # by row processing
        with open(self.path, 'r', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            csvRows = [row for row in reader]
        Utils.output(json.dumps(csvRows, indent=4), _print=False)
        # by columns processing
        csvCols = defaultdict(list) # each value in each column is appended to a list
        with open(self.path, 'r', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            for row in reader: # read a row as { column1: value1, column2: value2, ... }
                for (k, v) in row.items(): # for each column (k, v) pair
                    csvCols[k].append(v) # append the value into the appropriate list
        Utils.output(json.dumps(csvCols, indent=4), _print=False)
        csvContent = {
            "ByRows": csvRows,
            "ByCols": csvCols
        }
        _txt = Text(self.getRawText())
        csvContent["text-analysis"] = _txt.analysis()
        Utils.output(json.dumps(csvContent, indent=4), _print=False)
        return csvContent

    def extractXLSXinfo(self):

        # applies a string format for a value (checks for dict and list types):
        def formatValueToString(value):
            if isinstance(value, dict): # recursive
                formatted_value = {}
                for k, v in value.items():
                    formatted_value[k] = formatValueToString(v)
            elif type(value) is list: # recursive
                formatted_value = []
                for i in value:
                    formatted_value.append(formatValueToString(i))
            elif type(value) is datetime.datetime:
                formatted_value = f"{value.strftime('%Y-%m-%d %H:%M:%S')}"
            else:
                formatted_value = str(value)
            return formatted_value

        xlsxMetadata, xlsxBookDict = {}, {}
        book_dict = pe.get_book_dict(file_name=self.path, name_columns_by_row=1) # OrderedDict type
        # full update (without applying any format): xlsxBookDict.update(book_dict.items())
        for key, value in book_dict.items(): # applies a string format for all the values: specific case for datetime types.
            xlsxBookDict[key] = formatValueToString(value)
        Utils.output(json.dumps(xlsxBookDict, indent=4), _print=False)
        xlsxMetadata["workbooks"] = xlsxBookDict
        _txt = Text( str(xlsxBookDict) )
        xlsxMetadata["text-analysis"] = _txt.analysis()
        Utils.output(json.dumps(xlsxMetadata, indent=4), _print=False)
        return xlsxMetadata


    def isTextBased(self):
        return bool(re.search(self.extension.upper(), Utils._config["Text-Based-File-Extensions"])) 


    def getRawText(self, _encoding="utf-8"):
        txt, text_read_ok = "", False
        if (self.isTextBased()):
            '''
            Common character encoding issues:
            
            * UnicodeDecodeError: 'utf-8' codec can't decode byte 0xa9/0x85/0x92/0x96/0x97 in position xxx: invalid start byte
            Reason: "0x92 is a smart quote(’) of Windows-1252. It simply doesn't exist in Unicode, therefore it can't be decoded."
            Solution: "use encoding='cp1252'"
            Reference: <https://stackoverflow.com/questions/29419322/unicodedecodeerror-utf8-codec-cant-decode-byte-0x92-in-position-377826-inva>
            
            '''
            try_again = True
            for attempt in range(2): # two attempts: [0, 1]
                try:
                    if (try_again): # control flow of attempts
                        with open(self.path, encoding=_encoding) as f: # mode="rt" (default)
                            txt = f.read().encode(_encoding).decode(_encoding)
                except Exception as e: # if any failure occurs: return a dictionary with the Exception details.
                    if ((attempt == 0) and\
                        ("'utf-8' codec can't decode byte 0x" in str(e))): # try using 'cp1252'
                        _encoding = 'cp1252' # another attempt with the Windows-1252 encoding.
                        try_again = True
                    else: # any other exception
                        raise e
                else:
                    text_read_ok = True
                    break # exit the loop with success
            else: # all attempts fail
                text_read_ok = False
        else:
            Utils.output(f"The file is not text-based ('{Utils._config['Text-Based-File-Extensions']}').", _error=True, _print=True)
        return txt


    def extractZIPinfo(self):
        ZIP_targetDir = (Utils.getFolder("Temp-Folder") + self.name + "-" + str(random.randint(1,1000))) # with a random integer
        if not os.path.isdir(os.path.dirname(ZIP_targetDir)):
            Utils.output(f"Cannot create zipfile because target does not exists: target='{ZIP_targetDir}", _error=True, _print=True)
        else: # create the zipfile
            with zipfile.ZipFile(self.path,"r") as zip_ref:
                zip_ref.extractall(ZIP_targetDir)
            Utils.output(f"ZIP temporal directory: [{ZIP_targetDir}]", _print=False)
        zipMetadata, n = {}, 0
        for dirName, _, fileList in os.walk(ZIP_targetDir):
            _folder = os.path.basename(dirName)
            Utils.output(f"* Folder: '{dirName}'\n* Number of files found: {len(fileList)}", _print=False)
            for _filename in fileList: # for each file in the ZIP file.
                fileMetadata = {}
                n += 1
                Utils.output(f'... Extracting Metadata for file #{n}: [{_filename}]', _print=False)
                filepath = dirName + Utils.FILE_PATH_SEPARATOR + _filename
                f = File(path=filepath, useCase=self.useCase, defaultDirAttributes={}, JSONmetadata={}) # the default Directory attributes are set to empty.
                fileMetadata["folder"] = _folder
                fileMetadata["filename"] = _filename
                fileMetadata["metadata"] = f.extractMetadata()
                fileMetadata["content"] = f.getRawText() if f.isTextBased() else f.extractContent()
                zipMetadata[f"file-{n}"] = fileMetadata
                Utils.output(json.dumps(fileMetadata, indent=4), _print=False)
        Utils.output(json.dumps(zipMetadata, indent=4), _print=False)
        return zipMetadata


    def extractTXTinfo(self, content=""):
        txtMetadata = {}
        txt = Text(self.getRawText() if (not content) else content)
        txtMetadata["text-analysis"] = txt.analysis()
        Utils.output(json.dumps(txtMetadata, indent=4), _print=False)
        return txtMetadata


    def extractRTFinfo(self):
        RTF_text = self.getRawText()
        '''
        bookmarks, s = [], ""
        # Tried this but it took more than 4 hours in just one document!
        # r'({\\insrsid\d+\s|{\\\*\\bkmk(start|end) _Toc\d+})((: |\s)?\w+(\w+|\s*|\w\W)+){\\\*\\bkmk(start|end) _Toc\d+}'
        for m in re.finditer(r'{\\\*\\bkmk(start|end) _Toc\d+}((: |\s)?\w+(\w+|\s*|\w\W)+){\\\*\\bkmk(start|end) _Toc\d+}', RTF_text):
            s = m.group(2).strip()
            if (s):
                bookmarks.append(s)
        RTF_content["RTF-Bookmarks"] = bookmarks
        '''
        plain_text = rtf_to_text(RTF_text) # uses *striprtf*
        RTF_content = {}
        RTF_content.update(self.extractTXTinfo(plain_text))
        return RTF_content


    def extract_ext_info(self):
        Utils.output(f"extract_ext_info: ({self.extension})", _print=False)
        ext = self.extension.upper()
        extension_mapping = Utils._config["Extension-Mappings"].get(ext, ext)
        # Get the method from 'self'.  Default to a lambda.
        method = getattr(self, ("extract" + extension_mapping + "info"), lambda: {})
        return method()


    def extract_associated_NHMRC_info(self):
        global NHMRC_Grants
        NHMRC_data_file = Utils._config["Associated-Metadata"]["NHMRC-Settings"]
        # Ignore if it's the file that has the associated metadata:
        if (self.path !=
            NHMRC_data_file["Path"] + NHMRC_data_file["Folder"] + Utils.FILE_PATH_SEPARATOR + NHMRC_data_file["File"]):
            __fnc__ = sys._getframe().f_code.co_name + '#' + str(sys._getframe().f_code.co_firstlineno) # function + number of 1st line the source code
            if (__fnc__ not in Utils.IMPORTED_PACKAGES): # not in the array?
                #if ( (__name__ == "__main__") or ((__name__ == "MEL") and (__package__ == 'MEL')) ):
                import MEL.NHMRC_Grants as NHMRC_Grants
                #else:
                #    from . import NHMRC_Grants
                Utils.IMPORTED_PACKAGES.append(__fnc__) # marked as already imported
            appID = self.useCase
            source = ""
            if (NHMRC_Grants.FoR_App_Data.C_SUBPATH_2019 in self.path):
                source = NHMRC_Grants.FoR_App_Data.C_SOURCE_2019
            else:
                source = NHMRC_Grants.FoR_App_Data.C_SOURCE_History
            Utils.output(f"{self.path}\n{NHMRC_Grants.FoR_App_Data.C_SUBPATH_2019}\n{(NHMRC_Grants.FoR_App_Data.C_SUBPATH_2019 in self.path)}\n{source}", _print=False)
            _metadata = NHMRC_Grants.FoR_App_Data.get(source, appID)
            if not (_metadata): # This is incorrect: for NHMRC, we should always have associated metadata!
                Utils.output(f"*** Associated Metadata was not found! ***", _error=True, _print=True)
            Utils.output(json.dumps(_metadata, indent=4), _print=False)
            return _metadata
        return {}


    def extract_associated_DoEE_Species_CAsAndRPs_info(self):
        global DoEE_Species
        DoEE_Species_data_file = Utils._config["Associated-Metadata"]["DoEE_Species_CAsAndRPs-Settings"]
        # Ignore if it's the file that has the associated metadata:
        if (self.path !=
            DoEE_Species_data_file["Path"] + DoEE_Species_data_file["Folder"] + Utils.FILE_PATH_SEPARATOR + DoEE_Species_data_file["File"]):
            __fnc__ = sys._getframe().f_code.co_name + '#' + str(sys._getframe().f_code.co_firstlineno) # function + number of 1st line the source code
            if (__fnc__ not in Utils.IMPORTED_PACKAGES): # not in the array?
                #if ( (__name__ == "__main__") or ((__name__ == "MEL") and (__package__ == 'MEL')) ):
                import MEL.DoEE_Species as DoEE_Species
                #else:
                #    from . import DoEE_Species
                Utils.IMPORTED_PACKAGES.append(__fnc__) # marked as already imported
            _metadata = DoEE_Species.CAsAndRPs.get(self.path, self.name)
            Utils.output(json.dumps(_metadata, indent=4), _print=False)
            return _metadata
        return {}


    def extract_associated_FutureSOILS_info(self):
        global FutureSOILS
        __fnc__ = sys._getframe().f_code.co_name + '#' + str(sys._getframe().f_code.co_firstlineno) # function + number of 1st line the source code
        if (__fnc__ not in Utils.IMPORTED_PACKAGES): # not in the array?
            #if ( (__name__ == "__main__") or ((__name__ == "MEL") and (__package__ == 'MEL')) ):
            import MEL.FutureSOILS_PPP_Engine as FutureSOILS
            #else:
            #    from . import DoEE_Species
            Utils.IMPORTED_PACKAGES.append(__fnc__) # marked as already imported
        _metadata = FutureSOILS.ProcessingEngine(self.dirname, self.name) # current processing file (directory folder and filename)
        Utils.output(json.dumps(_metadata, indent=4), _print=True)
        return _metadata


    def extract_associated___JSON___info(self):
        if (self.extension.upper() == "JSON"): # only for JSON files:
            strContent = self.getRawText()
            if (Utils.isJSON(strContent)):
                _metadata = json.loads(strContent)
                Utils.output(json.dumps(_metadata, indent=4), _print=False)
                return _metadata
        return {}


    # Extract "Associated-Metadata" from external domain-specific files such as XLS spreadsheets.
    def extract_associated_info(self):
        if (Dataset.hasAssociatedMetadata()):
            for name, method_id in Utils._config["Associated-Metadata"]["Mappings"].items():
                match = re.search(name, self.path) # looks if the file could have some associated info.
                if (match):
                    Utils.output(f"extract_associated_info: '{name}' ==> '{method_id}' ", _print=False)
                    # Get the method from 'self'.  Default to a lambda.
                    method = getattr(self, ("extract_associated_" + method_id + "_info"), lambda: {})
                    return method()
        return {}


    def extractContent(self):
        if (not(self._AreAttributesSet) and (self.JSONmetadata)): # the case were we already have a metadata structure: for the NLP/NER results.
            return {}
        dict_content, dict_attachment = {}, {}
        if not(self._AreAttributesSet):
            Utils.output(f"The file attributes have not been set yet.", _error=True, _print=True)
        else:
            with open(self.path, "rb") as f:
                content = base64.b64encode(f.read())
            dict_content["content_type"] = self.MIMEtype
            dict_content["data"] = content.decode('ascii')
            dict_attachment[self.name] = dict_content
            Utils.output(json.dumps(dict_attachment, indent=4), _print=False)
        return dict_attachment


# ==================================================================================================
class Directory:
    structure = {}
    encrypted_PDFs = [] # list of encrypted/password-secured PDFs
    _process = False # Was the process() function call?
    
    @staticmethod
    def scanAndBuildStructure(path):
        Directory.structure = {}
        with os.scandir(path) as dir_entries:
            for entry in dir_entries: # https://docs.python.org/3/library/os.html#os.DirEntry
                file = {}
                info = entry.stat()
                dm = datetime.datetime.fromtimestamp(info.st_mtime)
                dc = datetime.datetime.fromtimestamp(info.st_ctime)
                # Default General-Metadata for any file: file attributes.
                # This is based on the expected output from the "NLZN Metadata Extractor"
                file["FILENAME"] = f"{entry.name}"
                file["SEPARATOR"] = '\\'
                file["PARENT"] = path
                file["CANONICALPATH"] = entry.path
                file["ABSOLUTEPATH"] = entry.path
                file["FILE"] = f"{entry.is_file()}"
                file["DIRECTORY"] = f"{entry.is_dir()}"
                file["FILELENGTH"] = info.st_size
                if entry.is_file(): # gets extension and MIME type
                    f = File(path=entry.path, useCase="", defaultDirAttributes={}, JSONmetadata={})
                    file["EXTENSION"] = f.extension
                    file["TYPE"] = f.MIMEtype
                file["MODIFIED"] = f"{dm}"
                file["DATE"] = f"{dm.strftime('%Y%m%d')}"
                file["DATEPATTERN"] = "yyyyMMdd"
                file["TIME"] = f"{dm.strftime('%H%M%S')}"
                file["TIMEPATTERN"] = "HHmmSS"
                file["CREATED"] = f"{dc}"
                file["CREATED-DATE"] = f"{dc.strftime('%Y%m%d')}"
                file["CREATED-TIME"] = f"{dc.strftime('%H%M%S')}"
                Directory.structure[f"{entry.name}"] = file
        if not (Directory.structure):
            Utils.output(f"Path not found, empty, or unable to scanned!\npath={path}", _error=True, _print=True)
            return
        Utils.output(json.dumps(Directory.structure, indent=4), _print=False)


# ==================================================================================================
# Main process:
    @staticmethod
    def process(rootDir):
        # Only one time:
        import time
        from progress.bar import IncrementalBar
        
        # Evaluation for folders and files: skip or perform processing?
        def doProcessing(current, only, startFrom, _doProcessing):
            ''' *only___* overrides *startFrom___* if set '''
            if (len(only) > 0): # *only___* is set
                _doProcessing = True if (current in only) else False # (do not process anything else!)
            else: # *only* is not set ==> evaluates *startFrom*
                _doProcessing = True if not(_doProcessing) and (current == startFrom) else _doProcessing
            #Utils._log.info(f"\n_doProcessing='{_doProcessing}")
            return _doProcessing

        def getExtension(filename):
            return os.path.splitext(filename)[1].translate({ord(c): None for c in '.'})
        def checkExtensionLogic(key, filename, emptyValueResult):
            flag = Dataset._db[key].strip()
            if ((flag == "") or (flag == "*")): # (empty value)
                return emptyValueResult
            extension = getExtension(filename)
            #Utils._log.info(f"\nflag='{flag}'; extension='{extension}'; return={extension.upper() in flag.upper()}")
            return (extension.upper() in flag.upper())

        def checkReloadFilesOfExtension(filename):
            return checkExtensionLogic("Reload-Files-of-Extensions", filename, True) # (empty value) means "any file type will be reloaded"

        def checkOnlyFilesOfExtension(filename): # check the processing of only specific file formats (extension)
            return checkExtensionLogic("Only-Files-of-Extensions", filename, True) # (empty value) means "any extension"

        def checkSkipFilesOfExtension(filename):
            return checkExtensionLogic("Skip-Files-of-Extensions", filename, False) # (empty value) means "do not skip any" ==> "process all"

        processed = 0
        def addDoc(f): # adds the document on CouchDB
            nonlocal processed
            _, result, metadata = CouchDB.addDocument(f) # performs the metadata extraction...
            if ((result is not None) and (result)):
                if ("ok" in result):
                    if (result["ok"] == True):
                        Utils._log.info(f'...... #{i}/{n}: OK! | id=[{result["id"]}]\n')
                        processed += 1
                        return
                if ("id" in result):
                    Utils._log.info(f'...... #{i}/{n}: error! | id=[{result["id"]}]\n')
                Utils._log.info(json.dumps(result, indent=4))
            else:
                Utils._log.info(f'...... #{i}/{n}: error! | *CouchDB.addDocument()* returned an empty result!\n')
            return metadata # returns the extracted metadata with, possibly, the NER results.

        c, u, skipped, reloaded = 0, 0, 0, 0
        startFromUseCase  = Dataset._f["Start-From-UseCase" ].strip()
        startFromFilename = Dataset._f["Start-From-Filename"].strip()
        onlyUseCase  = Dataset._f["Only-UseCase" ].strip()
        onlyFilename = Dataset._f["Only-Filename"].strip()
        doFolderProcessing = False if (len(startFromUseCase)  > 0) and (startFromUseCase  != "*") else True
        doFileProcessing   = False if (len(startFromFilename) > 0) and (startFromFilename != "*") else True
        for dirName, subdirList, fileList in os.walk(rootDir):
            if (dirName != rootDir):
                m = len(subdirList)
                Utils._log.info(f"Directory: [{dirName}]")
                Utils._log.info(f"Number of sub-directories found in it: {m}\n")
                _useCase = os.path.basename(dirName)
                n = len(fileList)
                u += 1
                Utils._log.info(f"Use Case (current directory): '{_useCase}'")
                Utils._log.info(f"Overall Use-Case counter: #{u}")
                Utils._log.info(f"Number of documents found: {n}\n")
                doFolderProcessing = doProcessing(_useCase, onlyUseCase, startFromUseCase, doFolderProcessing)
                if not(doFolderProcessing):
                    Utils._log.info(f"...... skipping this use case ......")
                else: # process the current folder
                    Directory.scanAndBuildStructure(dirName) # current directory structure: Directory.structure
                    i = 1
                    barFile = IncrementalBar('Countdown', max = n)
                    for _filename in fileList:
                        c += 1
                        doFileProcessing = doProcessing(_filename, onlyFilename, startFromFilename, doFileProcessing)
                        if ( not(doFileProcessing) or \
                             not(checkReloadFilesOfExtension(_filename)) or \
                             not(checkOnlyFilesOfExtension  (_filename)) or \
                                (checkSkipFilesOfExtension  (_filename)) ): # checks extension for processing:
                            Utils._log.info(f"...... (#{c}) skipping file: [{_filename}]")
                            skipped += 1
                            continue # next iteration
                        # process the current file
                        Utils._log.info(f'...... Overall file counter: #{c}')
                        Utils._log.info(f'...... Processing #{i} out of {n} ({round(i/n*100, 2)}%): [{_filename}]')
                        # file attributes are set in the constructor:
                        # default Directory attributes for the specified filename.
                        f = File(path=(dirName + Utils.FILE_PATH_SEPARATOR + _filename), useCase=_useCase, defaultDirAttributes=Directory.structure[_filename], JSONmetadata={})
                        output = {}
                        if (f.isEmpty() or f.isProtected()):
                            Utils._log.info(f"...... (#{c}) skipping file: [{_filename}]")
                        else:
                            if (Dataset.couchDBenabled()):
                                q = CouchDB.findDocument(f) # looks for the document on CouchDB
                                if q["execution_stats"]["results_returned"] >= 1: # if exists...
                                    Utils._log.info(f'...... #{i}/{n}: OK! | Already on CouchDB | id=[{q["docs"][0]["_id"]}]\n')
                                    if Dataset.reloadFile(f.extension): # reload file?
                                        CouchDB.deleteDocument(q["docs"][0]["_id"], q["docs"][0]["_rev"])
                                        output = addDoc(f) # extracted metadata
                                        reloaded += 1
                                else:
                                    output = addDoc(f) # extracted metadata
                            Dataset.generateOutputFile(_filename, output if (output) else f, processNER=Dataset.NER_OnlyNER())
                        i += 1
                        barFile.next()
                        time.sleep(0.5)
                    barFile.finish()
                Utils._log.info(f"\n******************************************************************************")
        Utils._log.info(f"\nStats:\nFolders: #{u}\nFiles: #{c}\nSkipped: #{skipped}\nReloaded: #{reloaded}\nProcessed: #{processed}")
        Utils._log.info(f"\n\nList of encrypted/password-secured PDFs found during processing:")
        if (Directory.encrypted_PDFs):
            i = 1
            for _pdf in Directory.encrypted_PDFs:
                Utils._log.info(f"{i}: {_pdf}")
                i += 1
        else:
            Utils._log.info(f"...... ** No encrypted/password-secured PDFs were found **")
        Directory._process = True


# ==================================================================================================
class NER:
    config = Utils._config["NLP-NER"]
    db_px  = config["Database-Prefix"]
    models: list = []
    processing_block: list = []
    loaded: bool = False
    _processed_docs: list = []
    _last_doc: dict = {}
    _last_hash_doc: str = ""
    _last_output_filename_template:str = ""
    _last_result: dict = {}
    _last_result_summary:dict = {}
    _last_output: dict = {}


    @staticmethod
    def models():
        NER.loaded = False
        NER.models = []
        NER.processing_block = []
        block_names = []
        
        for block in NER.config["Processing-Blocks"]:
            if (bool(block["Enable"])):
                block_names.append(block["Name"])
                NER.processing_block.extend(block["$"])

        isModelEnabled = lambda m : (bool(NER.config["Models"][m]) and (m in NER.processing_block))
        if isModelEnabled("spacy_sm"): NER.models.append(_NLP_NER.NERUtils.spacy_sm_name)
        if isModelEnabled("spacy_md"): NER.models.append(_NLP_NER.NERUtils.spacy_md_name)
        if isModelEnabled("spacy_lg"): NER.models.append(_NLP_NER.NERUtils.spacy_lg_name)
        if isModelEnabled("stanford_class3"): NER.models.append(_NLP_NER.NERUtils.stanford_class3_name)
        if isModelEnabled("stanford_class4"): NER.models.append(_NLP_NER.NERUtils.stanford_class4_name)
        if isModelEnabled("stanford_class7"): NER.models.append(_NLP_NER.NERUtils.stanford_class7_name)
        if isModelEnabled("stanza"): NER.models.append(_NLP_NER.NERUtils.stanza_name)
        if isModelEnabled("nltk"  ): NER.models.append(_NLP_NER.NERUtils.nltk_name)
        if isModelEnabled("bert"  ): NER.models.append(_NLP_NER.NERUtils.bert_name)
        if isModelEnabled("flair" ): NER.models.append(_NLP_NER.NERUtils.flair_name)
        if isModelEnabled("flair_ontonotes"     ): NER.models.append(_NLP_NER.NERUtils.flair_ontonotes_name)
        if isModelEnabled("flair_fast"          ): NER.models.append(_NLP_NER.NERUtils.flair_fast_name)
        if isModelEnabled("flair_fast_ontonotes"): NER.models.append(_NLP_NER.NERUtils.flair_fast_ontonotes_name)
        if isModelEnabled("flair_pooled"        ): NER.models.append(_NLP_NER.NERUtils.flair_pooled_name)
        if isModelEnabled("deeppavlov_onto"     ): NER.models.append(_NLP_NER.NERUtils.deeppavlov_onto_name)
        if isModelEnabled("deeppavlov_onto_bert"): NER.models.append(_NLP_NER.NERUtils.deeppavlov_onto_bert_name)
        if isModelEnabled("deeppavlov_conll2003"): NER.models.append(_NLP_NER.NERUtils.deeppavlov_conll2003_name)
        if isModelEnabled("deeppavlov_conll2003_bert"): NER.models.append(_NLP_NER.NERUtils.deeppavlov_conll2003_bert_name)
        if isModelEnabled("allennlp"): NER.models.append(_NLP_NER.NERUtils.allennlp_ner_name)
        if isModelEnabled("allennlp_finegrained"): NER.models.append(_NLP_NER.NERUtils.allennlp_finegrained_ner_name)
        if isModelEnabled("polyglot"): NER.models.append(_NLP_NER.NERUtils.polyglot_name)

        Utils.output(f"NER.processing_blocks={block_names}\nNER.models={NER.models}", _print=True)
        return NER.models


    @staticmethod
    def load_models():
        if ((Dataset.NERenabled()) and (not NER.loaded)):
            NER.models()
            NLP_NER.load_models(NER.models)
            NER.loaded = True # only loads the models once


    @staticmethod
    def process(mdDocsStructure={}, models2process=None): # returns the structure: o["NLP-NER"]
        # if the document has been already processed -> skip it!
        NER._last_result = { }
        if ("NLP-NER"         in mdDocsStructure):  NER._last_result["NLP-NER"]         = mdDocsStructure["NLP-NER"]  # current results
        if ("NLP-NER-Summary" in mdDocsStructure):  NER._last_result["NLP-NER-Summary"] = mdDocsStructure["NLP-NER-Summary"]  # current results
        _doc: str = mdDocsStructure["General-Metadata"]["FILENAME"]
        if (_doc in NER._processed_docs):
            Utils.output(f"\n[TNNT WARNING]: The document [{_doc}] was previously processed.  Skipping it...\n", _print=True)
            return NER._last_result
        models2process = NER.models if (models2process is None) else models2process # If *None* -> use *default* model list
        # *models2process* might be an empty list: [] --> do nothing!
        if ((Dataset.NERenabled()) and (NER.loaded) and (models2process)): # We assume that the models have been loaded:
            Utils.output(f"Processing the following models: {models2process}", _print=False)
            NER._last_result["NLP-NER"] = NLP_NER.NER_main(models2process, mdDocsStructure)
            NER._last_result["NLP-NER-Summary"] = NER.postProcessing()
            NER._processed_docs.append(_doc)  # add the document in the list of "already processed" docs.
        return NER._last_result


    @staticmethod
    def output():
        if not(NER._last_result) or ("NLP-NER" not in NER._last_result) or ((not NER._last_result["NLP-NER"][0])): # empty results
            Utils.output(f"*NER._last_result* is empty or 'NLP-NER' dictionary not found.", _print=True)
            return
        results = NER._last_result["NLP-NER"][0]["doc-0"]
        if (NER._last_doc["General-Metadata"]["FILENAME"] == results["filename"]):
            # calculate hash for the "absolutepath":
            h = blake2s(digest_size=7)
            h.update(NER._last_doc["General-Metadata"]["ABSOLUTEPATH"].encode())
            NER._last_hash_doc = h.hexdigest()
            sufix = "--NER--(%%MODEL%%)--" + NER._last_hash_doc + ".json"
            file_path = Dataset._ner["Output"] 
            NER._last_output_filename_template = NER._last_doc["General-Metadata"]["FILENAME"] + sufix
            NER._last_output = { "NLP-NER": {} }
            for k, v in results.items():
                if (("_model" in k) or ("allennlp" in k)): # for each model
                    output = {}
                    output_filename = file_path + NER._last_output_filename_template.replace("%%MODEL%%", k)
                    generate = True if not(os.path.exists(output_filename)) else Dataset.NER_regenerateOutput()
                    if (generate):
                        output.update(NER._last_doc)  # document header
                        output["NLP-NER"] = { k : v } # NER model result
                        output["NLP-NER-Summary-(From-Last-Run)"] = NER._last_result_summary  # NER model result summary
                        with open(output_filename, 'w') as outfile:
                            outfile.write(json.dumps(output, indent=4))
                        # NER model result ==> *indexing* the output
                        NER._last_output["NLP-NER"].update({ k : v })
                        NER._last_output["NLP-NER"][k]["_output"] = output_filename[len(file_path):] # only filename
                    Utils._log.info(f"[{'new' if (generate) else 'old'}] {output_filename}")
            Utils.output(json.dumps(NER._last_output, indent=4), _print=False)


    @staticmethod
    def processFromToCouchDB_forOneDoc(input_doc): # *input_doc* has the full structure of the document to be processed
        output = {}

        def processingNER(current_NER_results):
            Utils.output(f"... processing the models on the document (hold on)...", _print=True)
            models2process = NER.models.copy() # copy the list of models to process from the *default* model list
            if (current_NER_results): # if there are some preliminary results
                for m, _ in current_NER_results.items(): # for each model "m" in the results
                    if (m in models2process):
                        models2process.remove(m) # the model "m" has already been processed -> remove it from the models to process.
            NER.process(input_doc, models2process) # full result is stored in *NER._last_result*
            NER.output() # generated output files are stored in *NER._last_output*
            # updating the output:
            output["NLP-NER$Output-Files"] = NER._last_output_filename_template
            output["NLP-NER"] = current_NER_results # init the output with current NER results
            if not(NER._last_result):
                Utils.output(f"*NER._last_result* is empty.", _print=False)
                return
            for k, v in NER._last_output["NLP-NER"].items():
                output["NLP-NER"].update({ k : v }) # update output to the latest results

        def addDoc(): # adds the NLP/NER results on CouchDB
            doc = File(JSONmetadata=output, path="", useCase="", defaultDirAttributes={})
            http_res, add_res, _ = CouchDB.addDocument(doc, CouchDB._NER_end_point)
            if ( ((add_res is not None) and (add_res)) and ((http_res is not None) and (http_res)) ):
                if ("ok" in add_res):
                    if (add_res["ok"] == True):
                        Utils._log.info(f'... OK! -- id=[{add_res["id"]}]\n')
                        return
                if ("id" in add_res):
                    Utils._log.info(f'... error! -- id=[{add_res["id"]}]\n')
                Utils._log.info(json.dumps(add_res, indent=4))
                if ((str(http_res) == "<Response [200]>") or (str(http_res) == "<Response [201]>")):
                    Utils.output(f'... OK! -- The request has succeeded (the doc. has been added)\n', _print=True)
                else:
                    Utils.output(f'... error! -- [{http_res}]\n', _error=True, _print=True)
                '''
                **HTTP Response Status Codes**
                ==============================
                + HTTP 409: "Conflict"
                indicates a request conflict with current state of the target resource. Conflicts are most likely to occur in response to a PUT request.
                + HTTP 413: "Payload Too Large"
                the request entity is larger than limits defined by server; the server might close the connection or return a Retry-After header field.
                '''
            else:
                Utils._log.info(f'... error! -- *CouchDB.addDocument()* returned an empty result!\n')


        if not(Dataset.couchDBenabled()):
            return
        # *input_doc* has the full structure of the document to be processed
        output["General-Metadata"] = input_doc["General-Metadata"]
        output["Use-Case$Folder" ] = input_doc["Use-Case$Folder" ]
        NER._last_doc = output
        Utils.output(json.dumps(output, indent=4), _print=False)
        if (Dataset.NER_storeOutputOnCouchDB()):
            f = File(JSONmetadata=output, path="", useCase="", defaultDirAttributes={})
            q = CouchDB.findDocument(f, Dataset._ner["Database"]) # looks for the document on CouchDB NLP-NER db.
            if q["execution_stats"]["results_returned"] >= 1: # if exists...
                Utils._log.info(f'... Already on CouchDB | id=[{q["docs"][0]["_id"]}]\n')
                d = CouchDB.getDocumentGivenID(q["docs"][0]["_id"], Dataset._ner["Database"]) # get current NER results from the dataset's NER DB
                processingNER(d["NLP-NER"]) # <-- execution of the NER task for the *input_doc* | update the output with current NER results: *d["NLP-NER"]*
                CouchDB.deleteDocument(q["docs"][0]["_id"], q["docs"][0]["_rev"], CouchDB._NER_end_point) # CouchDB endpoint (with NLP-NER db).
                addDoc()
            else:
                processingNER({}) # <-- execution of the NER task for the *input_doc* | first time (init output): *{}*
                addDoc()
        else:
            processingNER({}) # <-- execution of the NER task for the *input_doc* | taken as if it's running for the first time (init output): *{}*
        Dataset.generateOutputFile(output["General-Metadata"]["FILENAME"], output)
        Utils.output(json.dumps(output, indent=4), _print=False)


    @staticmethod
    def processFromAndToCouchDB_forDocSet(selector):
        if not(Dataset.couchDBenabled()):
            return
        global DB
        fields = ["_id", "_rev", "General-Metadata", "Use-Case$Folder", "Specific-Metadata"]
        Utils.output(f"* selector:\n{json.dumps(selector, indent=4)}\n* fields={fields}\n", _print=True)
        r = CouchDB.queryDocs(DB, selector, False, fields) # from the dataset's default DB
        n = 1
        for input_doc in r["docs"]:
            Utils._log.info(f"""\n#{n}: (id={input_doc["_id"]}) f={input_doc["General-Metadata"]["ABSOLUTEPATH"]}\n""")
            NER.processFromToCouchDB_forOneDoc(input_doc)
            n += 1
    
    
    @staticmethod
    def postProcessing():
        ##with open(Dataset._f["Output-Folder"] + "BP4 2001-02.pdf-MEL+NER_output.json") as _f: # mode="rt" (default)
        ##    NER._last_result = json.loads(_f.read())
        _delimit = "|/*-$#~|"
        NER._last_result_summary = {}
        _fineGrained = {}
        _totals = {}
        _num_docs = 0
        for doc_set in NER._last_result["NLP-NER"]:
            for (d, r) in doc_set.items(): # (document, results)
                Utils.output(f"* document: {d}", _print=False)
                _num_docs += 1
                _fineGrained[d] = {}
                _totals[d] = {}
                NER._last_result_summary[d] = {}
                for (m, o) in r.items():
                    if (m == "filename"):
                        continue
                    Utils.output(f"* model: {m}", _print=False)
                    for (c, list_ents) in o["_output"].items():
                        Utils.output(f"  * category: {c}", _print=False)
                        for e in list_ents:
                            Utils.output(f"    * entity: {e['entity']}", _print=False)
                            Utils.addItemInCountingList(_fineGrained[d], f"{e['entity']}{_delimit}{c}{_delimit}{m}")
                            Utils.addItemInCountingList(_totals[d], f"{e['entity']}")
        Utils.output(f"* Number of docs: {_num_docs}\n* Fine grained results: {json.dumps(_fineGrained, indent=4)}", _print=False)
        _item = ""
        _item_dict = {}
        for (d, r) in _fineGrained.items():
            for (k, v) in r.items():
                _item = k.split(_delimit)
                _item_dict = {
                    "model": _item[2],
                    "category": _item[1],
                    "count": v
                }
                if (_item[0] not in NER._last_result_summary[d]):
                    NER._last_result_summary[d][_item[0]] = []
                NER._last_result_summary[d][_item[0]].append(_item_dict)
        for (d, r) in _totals.items():
            for (k, v) in r.items():
                NER._last_result_summary[d][k].append({ "total": v })
        Utils.output(f"* Summary: {json.dumps(NER._last_result_summary, indent=4)}", _print=False)
        return NER._last_result_summary


    @staticmethod
    def Sampling_getSamplesNOTinNERdb(all_datasets=False):
        ds = [Dataset._id] if (not all_datasets) else []
        samples = NLP_NER.sampling_get_samples_not_in_ner_db(datasets=ds, models=NER.models)
        Utils.output(f"{json.dumps(samples, indent=4)}", _print=False)
        Dataset.writeOutputFile(f"MEL.NER.Sampling_getSamplesNOTinNERdb(all_datasets={all_datasets})", f"--log.json", samples)
        return samples
    
    
    @staticmethod
    def processSamplingSet():
        def displayProcessingDoc(n, header, input_doc, footer):
            desc =  f'extension={input_doc["General-Metadata"]["EXTENSION"].lower()}' if (type == "by_extension") else\
                    f'length={   input_doc["General-Metadata"]["FILELENGTH"]}'        if (type == "by_length") else ""
            Utils._log.info(
                f"...... #{n+1}: BEGIN\n" +\
                f"...... #{n+1}: {{{desc}}} (id={input_doc['_id']}) f={input_doc['General-Metadata']['ABSOLUTEPATH']}" if (header) else f"" +\
                f"...... #{n+1}: END.\n" if (footer) else f""
            )
        samples = NER.Sampling_getSamplesNOTinNERdb(True) # for all the datasets from the *sampling* configuration file.
        doc_count = 0
        doc_ds_count = 0
        desc = ""
        for (dataset, sampling_type) in samples.items():
            Utils._log.info(f"""\n* Dataset: *{dataset}*""")
            Dataset.load(dataset) # loads the setting variables of the dataset to be processed...
            doc_ds_count = 0
            for (type, sets) in sampling_type.items():
                Utils._log.info(f""".. Sample type: *{type}*; {len(sets)} sets found:""")
                set_count = 1
                for s in sets:
                    Utils._log.info(f""".... Set #{set_count}; {len(s)} documents found{':' if (len(s) > 0) else '. --' }""")
                    n = 0
                    for doc in s:
                        # try to get the document content with "_id"
                        input_doc = CouchDB.getDocumentGivenID(doc["_id"], db=Dataset._db["Database"])
                        if  (input_doc == { "error": "not_found", "reason": "missing" }):
                            displayProcessingDoc(n, True, doc, False)
                            Utils.output(f"The document was not found by *_id*; retrieving the document by *FILENAME* and *ABSOLUTEPATH*", _print=True)
                            NER.processFromAndToCouchDB_forDocSet(
                                {
                                    "General-Metadata": {
                                        "FILENAME":     doc["General-Metadata"]["FILENAME"],
                                        "ABSOLUTEPATH": doc["General-Metadata"]["ABSOLUTEPATH"]
                                    }
                                }
                            )
                            input_doc = doc # for displaying purposes below
                        else: # the document was retrieved by "_id"
                            displayProcessingDoc(n, True, input_doc, False)
                            NER.processFromToCouchDB_forOneDoc(input_doc)
                        displayProcessingDoc(n, False, input_doc, True)
                        n += 1
                        doc_ds_count += 1
                        doc_count += 1
                    if (n > 0):
                        Utils._log.info("\n")
                    set_count += 1
            Utils._log.info(f""". Documents processed: {doc_ds_count}\n""")
        Utils._log.info(f"""* Total number of documents processed: {doc_count}\n""")
        Dataset.loadDefault() # restores the *default* dataset variables


# Only imports NLP_NER module if process and dataset flags are on (from the default configuration file):
if (    ( (Dataset.NERenabled())  and (not NER.loaded) ) and\
        ( (__name__ == "__main__") or ((__name__ == "MEL") and (__package__ == 'MEL')) )\
    ):
    import NLP_NER
    import NLP_NER.EntityRecognition as _NLP_NER
    import NLP_NER.Sampling as Sampling
    init() # MEL.init() --> loads the configuration file.  Checks whether the NER parameters are activated.
    Utils._log.info(f"******************************************************************************")
    Utils._log.info(f"Loading NLP/NER models...")
    NER.load_models()
    Utils._log.info(f"NLP/NER models loaded!")
    Utils._log.info(f"******************************************************************************\n")


# ==================================================================================================
class VirtuosoUS: # Virtuoso Universal Server
    # import jaydebeapi # JDBC connector to VUS
    _end_point = Utils._config["Virtuoso"]["Host"]

    @staticmethod
    def execSPARQL(query):
        params = {
            "default-graph": "http://rsmdb01.server-domain/KG/Stage/0",
            "should-sponge": "soft",
            "query": query,
            "debug": "off",
            "timeout": "",
            "format": "application/json",
            "save": "display",
            "fname": ""
        }
        Utils.output(json.dumps(params, indent=4), _print=True)
        req = requests.post(VirtuosoUS._end_point, json=params)
        # req = requests.get(CouchDB._end_point + '/' + sparql)
        Utils.output(str(req), _print=True)
        result = json.loads(req.text)
        Utils.output(json.dumps(result, indent=4), _print=True)
        # HTTP Response:
        json_response = json.loads(req.content)
        Utils.output(f"""
            * HTTP Response Code = {req.status_code}
            * HTTP Headers = {req.headers}
            * JSON Response = {json.dumps(json_response, indent=4)}""", _print=True)

        if (req.status_code == 200): # 200 OK - Request completed successfully
            Utils.output(f"""*** Details below: {
                                  json.dumps(json_response, indent=4)}""", _print=False)
            return json_response
        else:
            # Possible HTTP Response:
            # Status Codes:    
            #    400 Bad Request - Invalid request
            #    401 Unauthorized - Read permission required
            #    500 Internal Server Error - Query execution error
            Utils.output(f"""The query execution failed:
                * HTTP Response Code = {req.status_code}
                * HTTP Headers = {req.headers}""", _error=True, _print=True)
            Utils.output(json.dumps(json_response, indent=4), _error=True, _print=True)
            return req

    #===========================================================================
    # @staticmethod
    # def Jena(query):
    #     jclass = "org.apache.jena.jdbc.JenaJDBC"
    #     conn_string = "jdbc:virtuoso://rsmdb01.server-domain:8890"
    #     conn_prop = {
    #         'user': "<USR>",
    #         'password': "<PWD>"
    #     }
    #     jarPath = "E:/_Maven/repository/org/apache/jena/jena-rdfconnection/3.13.1/jena-rdfconnection-3.13.1.jar"
    #     conn = jaydebeapi.connect(jclass, conn_string, conn_prop, jarPath)
    #     cursor = conn.cursor()
    #     query = """
    #     SELECT DISTINCT ?a
    #     WHERE  {
    #         ?a ?b ?b .
    #     }
    #     """
    #     cursor.execute(query)
    #     # do something with the results
    #     cursor.close()
    #     conn.close()
    #===========================================================================


# ==================================================================================================
# Main MEL functions:
# running time stamps:
dt_begin, dt_end = datetime.datetime.now(), datetime.datetime.now()
_displayed_Header = False
_displayed_Footer = False


def header():
    init() # initializes MEL with the default configuration file.
    global _displayed_Header
    if (_displayed_Header):
        return
    Utils._log.info(f"******************************************************************************")
    Utils._log.info(f"Project AGRIF | MEL: Metadata Extractor & Loader --> @dataset={Utils.DATASET_ID}")
    dt_begin = datetime.datetime.now()
    Utils.printStartTimeStamp(dt_begin)
    Utils._log.info(f"platform: [{platform.system()}]; __name__: [{__name__}]; __package__: [{__package__}]")
    Dataset.shortPrint()
    Utils._log.info(f"******************************************************************************")
    _displayed_Header = True

def process():
    header()
    rootDir = Utils.getFolder() # "Input-Folder"
    Directory.process(rootDir)
    footer()

def footer():
    global _displayed_Footer
    if (_displayed_Footer):
        return
    Utils._log.info(f"******************************************************************************")
    dt_end = datetime.datetime.now()
    delta = (dt_end - dt_begin).total_seconds()
    Utils.printEndTimeStamp(dt_end, delta)
    if Dataset.couchDBenabled() and Directory._process:
        CouchDB.checkDuplicateDocs(Dataset._db["Database"])
    _displayed_Footer = True


# ==================================================================================================
# Set of useful functions for supporting purposes.

class Support():

    @staticmethod
    def NHMRC_EmptyFolders():
        FILE_PATH = "E:/_temp/DepFin-Project/NHMRC/NHMRC_Files.txt"
        f = open(FILE_PATH, 'r')
        folder, size, l, n = "", "", 1, 1
        Utils._log.info("List of empty folders (no files found in them):\n")
        for line in f:
            folder = line.strip() if ("Directory of " in line.strip()) else folder
            size   = line.strip() if ((size == "") and (" bytes" in line.strip())) else ""
            if ("0 File(s)" in size): # The folder is empty:
                Utils._log.info(f"#{n}: {folder[46:]}")
                #Utils._log.info(f"#{n}: l={l}; folder=[{folder}]; size=[{size}]")
                folder, size = "", ""
                n += 1
            l += 1
        f.close()


# ==================================================================================================
# Set of functions for testing purposes.

class Test():

    @staticmethod
    def VUS_query():
        query = """SELECT *
        WHERE { ?s ?p ?o }"""
        VirtuosoUS.execSPARQL(query)

    @staticmethod
    def File(_useCase="", _input_file=""): # metadata extractor
        _useCase    = Dataset._f["Only-UseCase" ] if (not _useCase   ) else _useCase    # default value 
        _input_file = Dataset._f["Only-Filename"] if (not _input_file) else _input_file # default value
        #"025 MyEnvironment Inc.oft"
        # /*- With PPTX: #
        #"2017-7881 Assessment-Email-OE-Summary of ecology meeting-20180511_0.1.msg"
        # /*- With ERROR: #
        #"2018-8155-Assessment-Draft PD.msg" # PDF attachment: "CIB16605.03 R002 Rev 0.pdf"
        _d = Dataset._f["Test-Folder"] + _useCase 
        Directory.scanAndBuildStructure(_d)
        _f = File(path=(_d + Utils.FILE_PATH_SEPARATOR + _input_file), useCase=_useCase, defaultDirAttributes=Directory.structure[_input_file], JSONmetadata={})
        _m = _f.extractMetadata()  # if (Dataset["NLP-NER"]["Output-Handling"] == "COMPLETE") is set in the configuration file --> NER.process()
        if (Dataset.isViewFilePrintEnabled()):
            _f.print()
        Dataset.generateOutputFile(_input_file, _m, processNER=Dataset.NER_Complete())  # It might run NER.process() twice!
        return _f

    @staticmethod
    def CouchDB():
        global DB

        def countDocsXsubPath():
            SUBPATH = "\\\\ps"
            #"Project Grants"
            #"Mammals - (.)+"
            #"CAsAndRPs\\\\recovery_plans"
            n = CouchDB.countDocsXsubPath(DB, SUBPATH)
            Utils.output(f"db={{{DB}}}: {SUBPATH} ==> {n} documents were found!", _print=True)

        def queryDocs():
            ''' MANGO Query examples:
            # Documents without Associated-Metadata @ /doee-species[CA_ps]:
            {
               "selector": {
                  "General-Metadata": {
                     "ABSOLUTEPATH": {
                        "$regex": "(.*)(conservation_advices)(.*)"
                     }
                  },
                  "Use-Case$Folder": "ps",
                  "Associated-Metadata": []
               }
            }
            # ("96-conservation-advice.pdf | 99-conservation-advice.pdf")
            {
               "selector": {
                  "General-Metadata": {
                     "FILENAME": {
                        "$regex": "^(9)[69]{1}(-conservation-advice\\.pdf)$"
                     }
                  }
               }
            }
            # NHMRC:
            {
               "selector": {
                  "General-Metadata": {
                     "ABSOLUTEPATH": {
                        "$regex": "(.*)(Missing Ideas Grants proposals)(.*)"
                     }
                  }
               }
            }
            # DoF - Procurements:
            {
               "selector": {
                  "Specific-Metadata": {}
               }
            }
            # DoF-PBS-data:
            {
                "selector": {
                    "General-Metadata": {
                        "$or": [
                            { "EXTENSION": "pdf" },
                            { "EXTENSION": "PDF" }
                        ]
                    }
                }
            }
            {
                "selector": {
                    "General-Metadata": { "EXTENSION": "htm" }
                }
            }
            # ###
            selector = {
                "Specific-Metadata": {}
            }
            '''
            selector = {
                "General-Metadata": { "EXTENSION": "rtf" }
            }
            fields = ["_id", "General-Metadata", "Use-Case$Folder", "Specific-Metadata"]
            #fields = ["Use-Case$Folder", "General-Metadata"]
            Utils.output(f"* selector:\n{json.dumps(selector, indent=4)}\n* fields={fields}\n", _print=True)
            r = CouchDB.queryDocs(DB, selector, False, fields)
            #Utils._log.info(f"db={{{DB}}}: ==> {r} documents were found!", _print=True)
            #Utils._log.info(f"db={{{DB}}}:\n{json.dumps(r, indent=4)}", _print=False)
            n = 1
            for x in r["docs"]:
                Utils._log.info(f"""#{n}: f={x["General-Metadata"]["ABSOLUTEPATH"]}\nid={x["_id"]}""")
                #Utils._log.info(f"""{n}. {x["Use-Case$Folder"]}\\{x["General-Metadata"]["ABSOLUTEPATH"]}""")
                n += 1

        def deleteDoc(_f):
            q = CouchDB.findDocument(_f)
            CouchDB.deleteDocument(q["docs"][0]["_id"], q["docs"][0]["_rev"])
        

        def NHMRC_pruneMetadata():
            x, n, i = {}, 0, 0
            def addDoc(x): # adds the pruned document on CouchDB
                f = File(JSONmetadata=x, path="", useCase="", defaultDirAttributes={})
                http_res, add_res, metadata = CouchDB.addDocument(f)
                if ((add_res is not None) and (add_res)):
                    if ("ok" in add_res):
                        if (add_res["ok"] == True):
                            Utils._log.info(f'... #{i}/{n}: OK! | id=[{add_res["id"]}]\n')
                            return
                    if ("id" in add_res):
                        Utils._log.info(f'... #{i}/{n}: error! | id=[{add_res["id"]}]\n')
                    Utils._log.info(json.dumps(add_res, indent=4))
                else:
                    Utils._log.info(f'... #{i}/{n}: error! | *CouchDB.addDocument()* returned an empty result!\n')
            Dataset.load("NHMRC-Grants") # NHMRC dataset
            Utils._log.info(f"** NHMRC - Pruning metadata task **")
            selector = { } # all documents
            fields = ["_id", "_rev", "General-Metadata", "Use-Case$Folder", "Specific-Metadata", "Associated-Metadata", "_attachments" ]
            r = CouchDB.queryDocs(DB, selector, False, fields)
            n = len(r['docs'])
            Utils._log.info(f"db={{{DB}}}: ==> {n} documents were found!")
            i = 1
            for x in r["docs"]:
                Utils._log.info(f"""#{i}: f={x["General-Metadata"]["ABSOLUTEPATH"]}\nid={x["_id"]}\nrev={x["_rev"]}""")
                # Pruning task: remove the following information from the JSON structure:
                if ("text-analysis" in x["Specific-Metadata"]): # (a) /"Specific-Metadata"/"text-analysis": it has the extracted content and text analysis.
                    x["Specific-Metadata"]["text-analysis"] = { }
                # (b) /"Associated-Metadata"/*, except /"Associated-Metadata"/"ANDS_Grant_Data": all the information that came from the Excel files (structure data).
                if ("ANDS_Grant_Data" in x["Associated-Metadata"]):
                    x["Associated-Metadata"] = { "ANDS_Grant_Data": x["Associated-Metadata"]["ANDS_Grant_Data"] }
                else:
                    x["Associated-Metadata"] = { }
                x["_attachments"] = { } # (c) /"_attachments": the document itself.
                CouchDB.deleteDocument(x["_id"], x["_rev"])
                x.pop('_id', None)
                x.pop('_rev', None)
                addDoc(x) # add new document
                i += 1
            Utils._log.info(f"* Documents processed: {i-1}/{n}")


        def searchContent():
            pattern = re.compile("ACKNOWLEDGEMENTS")
            #results = CouchDB.searchFirstMatch(_db=DB, _re_path="", _re_filename="2005", _list_extensions=[ "doc" ], _limit=5, _re_txt_search_pattern=pattern)
            results = CouchDB.searchAllMatches(_db=DB, _re_path="NHMRC_True_Positives", _re_filename="16451079", _list_extensions=[ "pdf" ], _limit=5, _re_txt_search_pattern=pattern)
            Utils._log.info(f"\n\n** RESULTS **")
            for doc in results:
                Utils._log.info(f"""\nf={doc["General-Metadata"]["ABSOLUTEPATH"]}\nr={doc["$RESULT$"]}""")
                if ( isinstance(doc["$RESULT$"], list) ):
                    n = 1
                    for r in doc["$RESULT$"]:
                        Utils._log.info(f"match #{n}: {r.group(0)}")
                        n += 1
                else:
                    Utils._log.info(f"""\nmatch={doc["$RESULT$"].group(0)}""")
            Dataset.writeOutputFile("testing", "-search_output.txt", results)
        
        def getChunkFromContent():
            _BEGIN = "Output cost attribution"
            _END = " basis\."
            #results = CouchDB.getFirstChunk(_db=DB, _re_path="", _re_filename="2005", _list_extensions=[ "doc" ], _limit=5, \
            #                                _re_txt_search_BEGIN_pattern=_BEGIN, \
            #                                _re_txt_search_END_pattern=_END)
            results = CouchDB.getAllChunks(_db=DB, _re_path="", _re_filename="2005", _list_extensions=[ "doc" ], _limit=5, \
                                            _re_txt_search_BEGIN_pattern=_BEGIN, \
                                            _re_txt_search_END_pattern=_END)
            Utils._log.info(f"\n\n** RESULTS **\n")
            NER_chunk_result = {}
            ss = "" # only one chunk
            c, d = 0, 0
            for doc in results:
                d += 1
                Utils._log.info(f"""Doc. #{d}, f={doc["General-Metadata"]["ABSOLUTEPATH"]}\nr={doc["$RESULT$"]}""")
                chunk = ""
                ss = "" # only one chunk
                if ( isinstance(doc["$RESULT$"], list) ):
                    n = 1
                    for r in doc["$RESULT$"]:
                        chunk += (("\n|***|\n") if (n > 1) else "") + r.group(0)
                        Utils._log.info(f"chunk #{n}:\n******\n{r.group(0)}\n******\n\n")
                        n += 1
                        c += 1
                    if (n-1 > 1):
                        ss = "s" # many
                else:
                    chunk = doc["$RESULT$"].group(0)
                    Utils._log.info(f"""\nchunk={chunk}\n""")
                NER_chunk_result = {
                    "General-Metadata" : doc["General-Metadata"],
                    "Specific-Metadata": {
                        "text-analysis": {
                            "clean-text" : chunk # chunk(s) to analyze
                        }
                    },
                    "$PATTERN$": doc["$PATTERN$"] # just for reference
                }
                NER.process(NER_chunk_result) # NER processing
                Dataset.generateOutputFile(doc["General-Metadata"]["FILENAME"] + f"-chunk{ss}", NER_chunk_result, processNER=False)
                Utils.output(f"{json.dumps(NER_chunk_result, indent=4)}\n", _print=True)
            ss = "" if (c == 1) else "s"
            Dataset.writeOutputFile("testing", f"-chunk{ss}_output.txt", results)

        #_f = Test.File()
        #CouchDB.addDocument(_f)
        #CouchDB.findDocument(_f)
        #deleteDoc(_f)
        
        NHMRC_pruneMetadata()
        
        #TEST_DOC_ID = "1b195dc7f6535d739616cb836a007097"
        #CouchDB.getDocumentGivenID(TEST_DOC_ID)
        
        #countDocsXsubPath()
        #queryDocs()
        
        #searchContent()
        #getChunkFromContent()


    @staticmethod
    def NER_process(filename, db=Dataset._db["Database"]): # NER processing
        #filename = "BP4 2001-02.pdf" # "/dof-pbs-data"
        # "2017-7881 Proposed Approval-PD-Att E-Decision notice.docm" # "/doee-assessments"
        # "2019-20 PBS - Defence Portfolio\\.pdf" # "/dof-pbs-data"
        r = CouchDB.getDocumentGivenFilename(filename, db)
        NER_result = NER.process(r) # NER processing
        Dataset.generateOutputFile(filename, r, processNER=False)
        Utils.output(f"{json.dumps(NER_result, indent=4)}", _print=True)

    @staticmethod
    def NER_sampling(datasets=None):
        datasets = [ Utils.DATASET_ID ] if (datasets is None) else datasets # default value
        NLP_NER.sampling_main(NER.models, datasets)


# ==================================================================================================
def main():
    header()
    #process()
    '''
    NER.processFromAndToCouchDB_forDocSet(
        {
            "General-Metadata": { "FILENAME": "2002-2003_DPRS_PBS_01_pref.htm" }
        }
    )
    '''
    #NER.processSamplingSet()
    #Test.File()
    #Test.CouchDB()
    #Test.NER_process("Gymnobelideus_leadbeateri_stats - 2017.xlsx")
        # "Gymnobelideus_leadbeateri_stats - 2017.xlsx" # "/doee-species"
        # "BP4 2001-02.pdf" # "/dof-pbs-data"
        # "2017-7881 Proposed Approval-PD-Att E-Decision notice.docm" # "/doee-assessments"
        # "2019-20 PBS - Defence Portfolio\\.pdf" # "/dof-pbs-data"
    #Test.NER_sampling()
    #Test.NER_sampling_get_samples_not_in_ner_db(True)
    footer()


# ==================================================================================================
# main
if (__name__ == "__main__"):
    main()

# ==================================================================================================
